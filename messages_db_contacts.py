#!/usr/bin/env python3
"""
Export an iMessage SQLite database (chat.db or sms.db) into a phone-style HTML conversation with contact name resolution.

Supports both database formats:
  - chat.db (from Mac or iTunes backup)
  - sms.db (from iPhone direct extraction)

Usage examples:
  # Basic with contact lookup
  python messages_db_contacts.py --db "chat.db" --addressbook "AddressBook-v22.abcddb" --out "conversation.html"

  # Using iPhone database
  python messages_db_contacts.py --db "sms.db" --addressbook "AddressBook-v22.abcddb" --out "conversation.html"

  # Wider phone look
  python messages_db_contacts.py --db "chat.db" --addressbook "AddressBook-v22.abcddb" --out "conversation.html" --width 420

  # Filter by chat and convert timezone
  python messages_db_contacts.py --db "chat.db" --addressbook "AddressBook-v22.abcddb" --out "conversation.html" --chat-id 1 --tz-from UTC --tz-to America/Chicago
"""

from __future__ import annotations
import argparse
import sys
import sqlite3
from pathlib import Path
from datetime import datetime, timedelta
from zoneinfo import ZoneInfo
import html
import math
import tempfile
import base64
from io import BytesIO
from typing import List, Tuple, Dict, Optional, Callable
import hashlib
import os
import threading
import json
import statistics

try:
    from tqdm import tqdm as _tqdm_original
    # Wrap tqdm to handle cases where stdout is None (GUI mode)
    def tqdm(iterable=None, **kwargs):
        # Check if stdout is available (None in GUI mode on Windows)
        if sys.stdout is None:
            # Return iterable if provided, otherwise return a dummy context manager
            if iterable is not None:
                return iterable
            else:
                # Return a dummy context manager for manual progress bars
                class DummyTqdm:
                    def __enter__(self):
                        return self
                    def __exit__(self, *args):
                        pass
                    def update(self, n=1):
                        pass
                return DummyTqdm()
        try:
            if iterable is not None:
                return _tqdm_original(iterable, **kwargs)
            else:
                return _tqdm_original(**kwargs)
        except:
            # If tqdm fails for any reason, return appropriate fallback
            if iterable is not None:
                return iterable
            else:
                class DummyTqdm:
                    def __enter__(self):
                        return self
                    def __exit__(self, *args):
                        pass
                    def update(self, n=1):
                        pass
                return DummyTqdm()
except ImportError:
    print("Warning: tqdm not found. Progress bars will be disabled.")
    print("Install tqdm for progress indicators: pip install tqdm")
    # Fallback: create a dummy tqdm that just returns the iterable
    def tqdm(iterable=None, **kwargs):
        if iterable is not None:
            return iterable
        else:
            class DummyTqdm:
                def __enter__(self):
                    return self
                def __exit__(self, *args):
                    pass
                def update(self, n=1):
                    pass
            return DummyTqdm()

ACTIVITY_FN_DEFAULT = None  # e.g., "my_participant_activity_pairs"

# Note: Print and tqdm wrappers above only activate when stdout is None (compiled EXE mode)
# In normal Python execution, they pass through to the original functions

# Apple's epoch starts on January 1, 2001
APPLE_EPOCH = datetime(2001, 1, 1)

def apple_timestamp_to_datetime(apple_timestamp):
    """Convert Apple timestamp to datetime object"""
    if apple_timestamp is None:
        return None
    # Apple timestamp is in nanoseconds since 2001-01-01
    seconds = apple_timestamp / 1000000000.0
    return APPLE_EPOCH + timedelta(seconds=seconds)

def resolve_attachment_path(attachment_path: str, attachments_folder: Optional[Path]) -> Optional[Path]:
    """
    Resolve an attachment path from the database to an actual file path.

    The database stores paths like ~/Library/Messages/Attachments/...
    This function tries to find the actual file by:
    1. Checking if the path exists as-is (expanding ~)
    2. If attachments_folder is provided, searching within it

    Returns the resolved Path if file exists, None otherwise.
    """
    if not attachment_path:
        return None

    # Try the path as-is with ~ expansion
    try:
        full_path = Path(attachment_path).expanduser()
        if full_path.exists() and full_path.is_file():
            return full_path
    except:
        pass

    # If attachments folder provided, try to find the file there
    if attachments_folder:
        try:
            # Extract just the filename from the database path
            filename = Path(attachment_path).name

            # Search for the file in the attachments folder
            attachments_root = Path(attachments_folder)
            if attachments_root.exists():
                # Try direct match first
                direct_match = attachments_root / filename
                if direct_match.exists() and direct_match.is_file():
                    return direct_match

                # Search recursively for the filename
                for file_path in attachments_root.rglob(filename):
                    if file_path.is_file():
                        return file_path
        except:
            pass

    return None

def detect_file_type_by_magic_bytes(file_path: Path) -> tuple[Optional[str], Optional[str]]:
    """
    Detect file type by reading magic bytes (file signature).
    Returns (mime_type, extension) or (None, None) if unrecognized.

    Common file signatures:
    - PNG: 89 50 4E 47 0D 0A 1A 0A
    - JPEG: FF D8 FF
    - GIF: 47 49 46 38
    - MP4: ftyp (at offset 4)
    - MOV: ftyp (at offset 4)
    - HEIC: ftyp (at offset 4) with heic/mif1
    """
    if not file_path or not file_path.exists():
        return None, None

    try:
        with open(file_path, 'rb') as f:
            header = f.read(12)  # Read first 12 bytes

        if len(header) < 4:
            return None, None

        # PNG: 89 50 4E 47
        if header[:4] == b'\x89PNG':
            return 'image/png', '.png'

        # JPEG: FF D8 FF
        if header[:3] == b'\xff\xd8\xff':
            return 'image/jpeg', '.jpg'

        # GIF: GIF8
        if header[:4] == b'GIF8':
            return 'image/gif', '.gif'

        # Check for MP4/MOV/HEIC (all use ftyp at offset 4)
        if len(header) >= 12 and header[4:8] == b'ftyp':
            ftype = header[8:12]
            # HEIC formats
            if ftype in [b'heic', b'heix', b'hevc', b'hevx', b'mif1']:
                return 'image/heic', '.heic'
            # MP4 formats
            elif ftype in [b'isom', b'iso2', b'mp41', b'mp42']:
                return 'video/mp4', '.mp4'
            # MOV/QuickTime
            elif ftype in [b'qt  ', b'M4V ', b'M4A ']:
                return 'video/quicktime', '.mov'

        # BMP: BM
        if header[:2] == b'BM':
            return 'image/bmp', '.bmp'

        # TIFF: II or MM
        if header[:2] in [b'II', b'MM']:
            return 'image/tiff', '.tiff'

        # WebP: RIFF....WEBP
        if header[:4] == b'RIFF' and len(header) >= 12 and header[8:12] == b'WEBP':
            return 'image/webp', '.webp'

        return None, None

    except Exception:
        return None, None

def embed_attachment_in_html(file_path: Path, mime_type: Optional[str]) -> Optional[str]:
    """
    Create HTML to embed an attachment.
    - Images: embedded as base64 data URIs (HEIC converted to JPEG)
    - Videos: embedded with HTML5 video tag
    - Other files: provide download link with file info

    Returns HTML string or None if file can't be read.
    """
    if not file_path or not file_path.exists():
        return None

    try:
        # Determine mime type if not provided
        extension = file_path.suffix.lower()

        # If no extension or unknown extension (like .pluginPayloadAttachment), try to detect by magic bytes
        if not extension or extension == '' or extension in ['.pluginpayloadattachment']:
            detected_mime, detected_ext = detect_file_type_by_magic_bytes(file_path)
            if detected_mime:
                mime_type = detected_mime
                extension = detected_ext

        if not mime_type:
            mime_map = {
                '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                '.png': 'image/png', '.gif': 'image/gif',
                '.heic': 'image/heic', '.heif': 'image/heif',
                '.mp4': 'video/mp4', '.mov': 'video/quicktime',
                '.m4v': 'video/x-m4v'
            }
            mime_type = mime_map.get(extension, 'application/octet-stream')

        file_size = file_path.stat().st_size

        # Handle images - embed as base64
        if mime_type and mime_type.startswith('image/'):
            try:
                # Check if it's a HEIC/HEIF file that needs conversion
                if extension in ['.heic', '.heif']:
                    try:
                        # Try to import pillow_heif for HEIC support
                        from PIL import Image
                        import pillow_heif

                        # Register HEIF opener with Pillow
                        pillow_heif.register_heif_opener()

                        # Open and convert HEIC to JPEG
                        img = Image.open(file_path)

                        # Convert to RGB if necessary (HEIC can have alpha channel)
                        if img.mode in ('RGBA', 'LA', 'P'):
                            background = Image.new('RGB', img.size, (255, 255, 255))
                            if img.mode == 'P':
                                img = img.convert('RGBA')
                            background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                            img = background
                        elif img.mode != 'RGB':
                            img = img.convert('RGB')

                        # Save as JPEG to BytesIO
                        output = BytesIO()
                        img.save(output, format='JPEG', quality=85)
                        file_data = output.getvalue()
                        mime_type = 'image/jpeg'

                    except ImportError:
                        # pillow_heif not available, try alternative method
                        try:
                            from PIL import Image
                            # Some systems can open HEIC with just Pillow
                            img = Image.open(file_path)
                            if img.mode != 'RGB':
                                img = img.convert('RGB')
                            output = BytesIO()
                            img.save(output, format='JPEG', quality=85)
                            file_data = output.getvalue()
                            mime_type = 'image/jpeg'
                        except Exception:
                            return f'<div style="color: #999; font-style: italic;">HEIC image format not supported. Install pillow-heif: pip install pillow-heif</div>'
                    except Exception as e:
                        return f'<div style="color: #999; font-style: italic;">Error converting HEIC image: {str(e)}</div>'
                else:
                    # Standard image formats
                    with open(file_path, 'rb') as f:
                        file_data = f.read()

                # Limit image size to 5MB for embedding
                if len(file_data) > 5 * 1024 * 1024:
                    return f'<div style="color: #999; font-style: italic;">Image too large to embed ({file_size:,} bytes)</div>'

                base64_data = base64.b64encode(file_data).decode('utf-8')
                return f'<img src="data:{mime_type};base64,{base64_data}" style="max-width: 100%; height: auto; border-radius: 10px; margin: 5px 0;" />'
            except Exception as e:
                return f'<div style="color: #999; font-style: italic;">Error loading image: {str(e)}</div>'

        # Handle videos - embed with video tag
        elif mime_type and mime_type.startswith('video/'):
            try:
                # Note: MOV/QuickTime files may not play in all browsers due to codec support
                # For best compatibility, videos should be in MP4 format with H.264 codec

                # For large videos or MOV files, provide a download link instead of embedding
                if file_size > 10 * 1024 * 1024 or extension in ['.mov']:
                    size_mb = file_size / (1024 * 1024)
                    note = "Video too large to embed" if file_size > 10 * 1024 * 1024 else "MOV format may not play in browser"
                    # Create file:// URL for the video - converts to proper URL format
                    file_url = file_path.as_uri()
                    return f'''<div style="padding: 10px; background-color: rgba(0,0,0,0.05); border-radius: 5px; margin: 5px 0;">
                        <div style="color: #666; margin-bottom: 5px;">🎬 Video: {html.escape(file_path.name)}</div>
                        <div style="color: #999; font-size: 0.9em; font-style: italic;">{note} ({size_mb:.1f} MB)</div>
                        <div style="margin-top: 8px;">
                            <a href="{file_url}" style="display: inline-block; padding: 8px 16px; background-color: #007bff; color: white; text-decoration: none; border-radius: 5px; font-size: 0.9em;">▶️ Open Video</a>
                        </div>
                        <div style="color: #666; font-size: 0.75em; margin-top: 8px; word-break: break-all;">Path: {html.escape(str(file_path))}</div>
                    </div>'''

                # For smaller MP4 files, try to embed
                with open(file_path, 'rb') as f:
                    file_data = f.read()
                    base64_data = base64.b64encode(file_data).decode('utf-8')
                    return f'<video controls style="max-width: 100%; border-radius: 10px; margin: 5px 0;"><source src="data:{mime_type};base64,{base64_data}" type="{mime_type}">Your browser does not support the video tag.</video>'
            except Exception as e:
                return f'<div style="color: #999; font-style: italic;">Error loading video: {str(e)}</div>'

        # For other files, just show info
        else:
            return f'<div style="color: #999; font-style: italic;">File attachment: {file_path.name} ({file_size:,} bytes)</div>'

    except Exception as e:
        return f'<div style="color: #999; font-style: italic;">Error processing attachment: {str(e)}</div>'

def normalize_phone_number(phone_number: str) -> str:
    """Normalize a phone number for comparison by removing non-digits"""
    if not phone_number:
        return ""
    # Remove all non-digit characters
    normalized = ''.join(filter(str.isdigit, phone_number))
    # For US numbers, consider last 10 digits as the main identifier
    if len(normalized) >= 10:
        return normalized[-10:]
    return normalized

class ContactLookup:
    """Class to handle contact name lookup from AddressBook database"""

    def __init__(self, addressbook_path: Optional[str] = None):
        self.contacts = {}  # phone_number -> contact_name
        self.fallback_contacts = {}  # Manual fallback mapping
        self.fallback_usage_log = {}  # Track when fallback contacts are used: phone -> count
        self.failed_lookups = {}  # Track phone numbers that couldn't be resolved: phone -> count
        self.debug_mode = False  # Enable for detailed logging

        if addressbook_path and Path(addressbook_path).exists():
            self._load_contacts_from_db(addressbook_path)

        # Add manual fallback mappings for common numbers you might know
        # You can extend this or load from a CSV file
        self._load_fallback_contacts()

    def _load_contacts_from_db(self, addressbook_path: str):
        """Load contacts from AddressBook database (supports both Mac and iPhone formats)"""
        try:
            print(f"\n🔍 Loading contacts from: {addressbook_path}")
            conn = sqlite3.connect(addressbook_path)
            cursor = conn.cursor()

            # Detect database format by checking for table existence
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
            tables = [row[0] for row in cursor.fetchall()]
            print(f"📋 Found {len(tables)} tables in AddressBook database:")
            print(f"   Tables: {', '.join(tables[:10])}{' ...' if len(tables) > 10 else ''}")

            if 'ZABCDPHONENUMBER' in tables and 'ZABCDRECORD' in tables:
                # Mac/iTunes backup format (AddressBook-v22.abcddb)
                print("✓ Detected Mac/iTunes AddressBook format")
                self._load_mac_contacts(cursor)
            elif 'ABPerson' in tables and 'ABMultiValue' in tables:
                # iPhone format (AddressBook.sqlitedb)
                print("✓ Detected iPhone AddressBook format")
                self._load_iphone_contacts(cursor)
            else:
                print(f"❌ Warning: Unknown AddressBook database format")
                print(f"   Expected tables not found. Looking for:")
                print(f"   - Mac format: ZABCDPHONENUMBER, ZABCDRECORD")
                print(f"   - iPhone format: ABPerson, ABMultiValue")

            conn.close()
            print(f"✓ Loaded {len(self.contacts)} contacts from AddressBook")

            # Show sample contacts for debugging
            if self.contacts:
                print(f"📞 Sample contacts (first 5):")
                for i, (phone, name) in enumerate(list(self.contacts.items())[:5]):
                    print(f"   {phone} → {name}")

        except Exception as e:
            print(f"❌ Warning: Could not load contacts from AddressBook: {e}")
            import traceback
            traceback.print_exc()

    def _load_mac_contacts(self, cursor):
        """Load contacts from Mac/iTunes AddressBook format"""
        # Join phone numbers with contact records to get name mappings
        cursor.execute("""
            SELECT p.ZFULLNUMBER, p.ZOWNER, p.ZLABEL, r.ZFIRSTNAME, r.ZLASTNAME, r.Z_PK
            FROM ZABCDPHONENUMBER p
            JOIN ZABCDRECORD r ON p.ZOWNER = r.Z_PK
            WHERE p.ZFULLNUMBER IS NOT NULL
            AND (r.ZFIRSTNAME IS NOT NULL OR r.ZLASTNAME IS NOT NULL)
        """)

        phone_mappings = cursor.fetchall()

        # Map phone numbers to names
        for phone_number, owner_id, label, first_name, last_name, record_pk in phone_mappings:
            name_parts = []
            if first_name:
                name_parts.append(first_name.strip())
            if last_name:
                name_parts.append(last_name.strip())

            if name_parts:
                full_name = ' '.join(name_parts)
                normalized_phone = normalize_phone_number(phone_number)
                if normalized_phone:
                    self.contacts[normalized_phone] = full_name

    def _load_iphone_contacts(self, cursor):
        """Load contacts from iPhone AddressBook format"""
        # Property 3 is phone numbers in ABMultiValue
        cursor.execute("""
            SELECT mv.value, p.First, p.Last, p.DisplayName, p.ROWID
            FROM ABMultiValue mv
            JOIN ABPerson p ON mv.record_id = p.ROWID
            WHERE mv.property = 3
            AND mv.value IS NOT NULL
        """)

        phone_mappings = cursor.fetchall()

        # Map phone numbers to names
        for phone_number, first_name, last_name, display_name, record_id in phone_mappings:
            # Build name from available fields (prioritize DisplayName, then First/Last)
            name_parts = []

            if display_name and display_name.strip():
                full_name = display_name.strip()
            else:
                if first_name and first_name.strip():
                    name_parts.append(first_name.strip())
                if last_name and last_name.strip():
                    name_parts.append(last_name.strip())
                full_name = ' '.join(name_parts) if name_parts else None

            if full_name:
                normalized_phone = normalize_phone_number(phone_number)
                if normalized_phone:
                    self.contacts[normalized_phone] = full_name

    def _load_fallback_contacts(self):
        """Load manual fallback contact mappings"""
        # You can customize this or load from a file
        # Format: normalized_phone_number -> name
        fallback_file = Path("contacts_fallback.txt")

        if fallback_file.exists():
            try:
                with open(fallback_file, 'r') as f:
                    for line in f:
                        line = line.strip()
                        if line and not line.startswith('#') and '=' in line:
                            phone, name = line.split('=', 1)
                            normalized_phone = normalize_phone_number(phone.strip())
                            if normalized_phone:
                                self.fallback_contacts[normalized_phone] = name.strip()
                print(f"Loaded {len(self.fallback_contacts)} fallback contacts")
            except Exception as e:
                print(f"Warning: Could not load fallback contacts: {e}")
        else:
            # Create a sample fallback file
            sample_content = """# Contact fallback mappings
# Format: phone_number = Contact Name
# Example:
# +1234567890 = John Doe
# 5551234567 = Jane Smith
"""
            try:
                with open(fallback_file, 'w') as f:
                    f.write(sample_content)
                print(f"Created sample fallback contact file: {fallback_file}")
            except:
                pass

    def get_contact_name(self, phone_number: str, show_phone: bool = False, debug: bool = False) -> Optional[str]:
        """Get contact name for a phone number, optionally including the phone number"""
        if not phone_number:
            return None

        normalized_phone = normalize_phone_number(phone_number)
        if not normalized_phone:
            if debug:
                print(f"   ⚠️  Could not normalize phone: {phone_number}")
            return None

        contact_name = None
        match_method = None

        # Strategy 1: Exact match (normalized)
        if normalized_phone in self.contacts:
            contact_name = self.contacts[normalized_phone]
            match_method = "exact"
        elif normalized_phone in self.fallback_contacts:
            contact_name = self.fallback_contacts[normalized_phone]
            match_method = "fallback-exact"
            self.fallback_usage_log[normalized_phone] = self.fallback_usage_log.get(normalized_phone, 0) + 1
        else:
            # Strategy 2: Try last 10 digits (for different country code formats)
            all_contacts = {**self.contacts, **self.fallback_contacts}
            if len(normalized_phone) >= 10:
                last_10 = normalized_phone[-10:]
                for stored_phone, name in all_contacts.items():
                    if len(stored_phone) >= 10 and stored_phone[-10:] == last_10:
                        contact_name = name
                        match_method = "last-10-digits"
                        break

            # Strategy 3: Try last 7 digits (for local matching)
            if not contact_name and len(normalized_phone) >= 7:
                last_7 = normalized_phone[-7:]
                for stored_phone, name in all_contacts.items():
                    if len(stored_phone) >= 7 and stored_phone[-7:] == last_7:
                        contact_name = name
                        match_method = "last-7-digits"
                        break

        # Track failed lookups
        if not contact_name:
            self.failed_lookups[phone_number] = self.failed_lookups.get(phone_number, 0) + 1

        if debug or self.debug_mode:
            if contact_name:
                print(f"   ✓ {phone_number} → {contact_name} (via {match_method})")
            else:
                print(f"   ✗ {phone_number} (normalized: {normalized_phone}) - NO MATCH")
                if len(self.contacts) > 0:
                    print(f"      Checked against {len(self.contacts)} contacts")

        # Format the result based on options
        if contact_name:
            if show_phone:
                return f"{phone_number} ({contact_name})"
            else:
                return contact_name
        else:
            return None

    def write_diagnostic_report(self, output_path: str = "contact_lookup_diagnostics.txt"):
        """Write a diagnostic report showing contact lookup statistics"""
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write("="*80 + "\n")
                f.write("CONTACT LOOKUP DIAGNOSTIC REPORT\n")
                f.write("="*80 + "\n\n")

                f.write(f"Total contacts loaded: {len(self.contacts)}\n")
                f.write(f"Fallback contacts: {len(self.fallback_contacts)}\n")
                f.write(f"Failed lookups: {len(self.failed_lookups)}\n")
                f.write(f"Fallback contacts used: {sum(self.fallback_usage_log.values())}\n\n")

                if self.failed_lookups:
                    f.write("-"*80 + "\n")
                    f.write("FAILED LOOKUPS (phone numbers not found in AddressBook):\n")
                    f.write("-"*80 + "\n")
                    sorted_failed = sorted(self.failed_lookups.items(), key=lambda x: x[1], reverse=True)
                    for phone, count in sorted_failed[:50]:  # Top 50
                        normalized = normalize_phone_number(phone)
                        f.write(f"{phone:20s} (normalized: {normalized:15s}) - {count:4d} lookups\n")
                    f.write("\n")

                if self.contacts:
                    f.write("-"*80 + "\n")
                    f.write("LOADED CONTACTS (sample):\n")
                    f.write("-"*80 + "\n")
                    for i, (phone, name) in enumerate(list(self.contacts.items())[:50]):
                        f.write(f"{phone:15s} → {name}\n")
                    f.write("\n")

            print(f"📄 Diagnostic report written to: {output_path}")
        except Exception as e:
            print(f"Warning: Could not write diagnostic report: {e}")

def timezone_choices():
    try:
        from zoneinfo import available_timezones
    except ImportError:
        available_timezones = None

    zones = []
    if available_timezones:
        try:
            zones = sorted(available_timezones())
        except Exception:
            zones = []

    if not zones:
        zones = [
            "UTC",
            "America/New_York",
            "America/Chicago",
            "America/Denver",
            "America/Los_Angeles",
            "Europe/London",
            "Europe/Paris",
            "Asia/Tokyo",
            "Australia/Sydney",
        ]
    return zones

TZ_CHOICES = timezone_choices()
DEFAULT_DEST_TZ = "America/Chicago"

def parse_date_filter(date_str: str, tz_zone=None) -> Optional[datetime]:
    """Parse date string for filtering (YYYY-MM-DD or YYYY-MM-DD HH:MM:SS)"""
    if not date_str:
        return None

    try:
        # Try parsing with time first
        if ' ' in date_str:
            parsed_dt = datetime.strptime(date_str, "%Y-%m-%d %H:%M:%S")
        else:
            # Just date, set to start of day
            parsed_dt = datetime.strptime(date_str, "%Y-%m-%d")

        # Add timezone if provided
        if tz_zone:
            parsed_dt = parsed_dt.replace(tzinfo=tz_zone)

        return parsed_dt
    except ValueError:
        print(f"Warning: Invalid date format '{date_str}'. Use YYYY-MM-DD or YYYY-MM-DD HH:MM:SS")
        return None

def datetime_to_apple_timestamp(dt: datetime) -> int:
    """Convert datetime to Apple timestamp (nanoseconds since 2001-01-01)"""
    if dt.tzinfo is None:
        # Assume UTC if no timezone
        dt = dt.replace(tzinfo=ZoneInfo('UTC'))

    # Convert to UTC for Apple timestamp calculation
    dt_utc = dt.astimezone(ZoneInfo('UTC'))

    # Calculate seconds since Apple epoch (2001-01-01 UTC)
    delta = dt_utc - APPLE_EPOCH.replace(tzinfo=ZoneInfo('UTC'))
    # Convert to nanoseconds
    return int(delta.total_seconds() * 1000000000)

def calculate_file_hashes(file_path):
    """Calculate multiple hash values for forensic verification"""
    try:
        file_path = Path(file_path)
        if not file_path.exists():
            return None

        # Initialize hash objects for different algorithms
        md5_hash = hashlib.md5()
        sha1_hash = hashlib.sha1()
        sha256_hash = hashlib.sha256()

        # Get file size and modification time
        file_stats = file_path.stat()
        file_size = file_stats.st_size
        modification_time = datetime.fromtimestamp(file_stats.st_mtime).strftime("%Y-%m-%d %H:%M:%S")

        # Read file in chunks to handle large databases efficiently
        chunk_size = 65536  # 64KB chunks
        total_chunks = (file_size + chunk_size - 1) // chunk_size

        with open(file_path, 'rb') as f:
            with tqdm(total=total_chunks, desc=f"Hashing {file_path.name}", unit="chunk", leave=False) as pbar:
                while chunk := f.read(chunk_size):
                    md5_hash.update(chunk)
                    sha1_hash.update(chunk)
                    sha256_hash.update(chunk)
                    pbar.update(1)

        hash_info = {
            'file_path': str(file_path),
            'file_size': file_size,
            'modification_time': modification_time,
            'md5': md5_hash.hexdigest().upper(),
            'sha1': sha1_hash.hexdigest().upper(),
            'sha256': sha256_hash.hexdigest().upper()
        }

        return hash_info

    except Exception as exc:
        print(f"Error calculating file hashes: {exc}")
        return None

def get_hash_verification_info(db_path, addressbook_path=None, pre_hashes=None, post_hashes=None, addressbook_pre_hashes=None, addressbook_post_hashes=None):
    """Generate hash verification information for forensic documentation"""
    try:
        hash_verification = {
            'verification_time': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'database_hashes': {
                'pre_processing': pre_hashes,
                'post_processing': post_hashes,
                'integrity_verified': False
            },
            'addressbook_hashes': {
                'pre_processing': addressbook_pre_hashes,
                'post_processing': addressbook_post_hashes,
                'integrity_verified': False
            } if addressbook_pre_hashes or addressbook_post_hashes else None
        }

        # Check database integrity
        if pre_hashes and post_hashes:
            hash_verification['database_hashes']['integrity_verified'] = (
                pre_hashes['md5'] == post_hashes['md5'] and
                pre_hashes['sha1'] == post_hashes['sha1'] and
                pre_hashes['sha256'] == post_hashes['sha256']
            )

        # Check AddressBook integrity
        if addressbook_pre_hashes and addressbook_post_hashes:
            hash_verification['addressbook_hashes']['integrity_verified'] = (
                addressbook_pre_hashes['md5'] == addressbook_post_hashes['md5'] and
                addressbook_pre_hashes['sha1'] == addressbook_post_hashes['sha1'] and
                addressbook_pre_hashes['sha256'] == addressbook_post_hashes['sha256']
            )

        return hash_verification

    except Exception as exc:
        print(f"Error generating hash verification info: {exc}")
        return None

def get_participant_activity_info(db_path, chat_id=None, contact_lookup=None, tz_from="UTC", tz_to=None, date_from_apple=None, date_to_apple=None):
    """Analyze participant activity patterns and response times"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Build query conditions for date filtering
        query_conditions = []
        query_params = []

        if chat_id:
            query_conditions.append("cmj.chat_id = ?")
            query_params.append(chat_id)

        if date_from_apple:
            query_conditions.append("cmj.message_date >= ?")
            query_params.append(date_from_apple)

        if date_to_apple:
            query_conditions.append("cmj.message_date <= ?")
            query_params.append(date_to_apple)

        where_clause = ""
        if query_conditions:
            where_clause = "WHERE " + " AND ".join(query_conditions)

        # Get total message count (all messages, not just 1-on-1) for reporting
        # Use LEFT JOIN to include orphaned messages (forensically important)
        total_count_query = f"""
            SELECT COUNT(DISTINCT m.ROWID)
            FROM message m
            LEFT JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
            {where_clause}
        """
        cursor.execute(total_count_query, query_params)
        total_messages_all = cursor.fetchone()[0]

        # Build query for participant analysis with date filtering
        # Only include 1-on-1 chats (where chat has exactly 1 handle participant + Me = 2 total)
        query = f"""
            SELECT
                m.ROWID,
                m.text,
                m.date,
                m.is_from_me,
                h.id as handle_identifier,
                cmj.message_date,
                LAG(m.date) OVER (ORDER BY cmj.message_date) as prev_date,
                LAG(m.is_from_me) OVER (ORDER BY cmj.message_date) as prev_is_from_me,
                LAG(h.id) OVER (ORDER BY cmj.message_date) as prev_handle,
                c.chat_identifier,
                c.display_name
            FROM message m
            JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
            LEFT JOIN chat c ON cmj.chat_id = c.ROWID
            LEFT JOIN handle h ON m.handle_id = h.ROWID
            WHERE cmj.chat_id IN (
                SELECT chat_id
                FROM chat_handle_join
                GROUP BY chat_id
                HAVING COUNT(DISTINCT handle_id) = 1
            )
            {" AND " + " AND ".join(query_conditions) if query_conditions else ""}
            ORDER BY cmj.message_date ASC
        """
        cursor.execute(query, query_params)

        messages = cursor.fetchall()

        if not messages:
            conn.close()
            return None

        # Initialize tracking variables
        participant_stats = {}
        response_times = []
        total_messages_1on1 = len(messages)  # Messages in 1-on-1 chats only

        # Process each message
        for msg in messages:
            rowid, text, date, is_from_me, handle_identifier, message_date, prev_date, prev_is_from_me, prev_handle, chat_identifier, display_name = msg

            # Determine participant identifier
            # For forensic accuracy, always use the actual person's handle/name, never group chat names
            if is_from_me:
                participant = "Me"
            else:
                # If handle_identifier is NULL, look up the chat's participant handle
                if not handle_identifier:
                    # Since we filtered to 1-on-1 chats only, there should be exactly one participant handle
                    temp_query = """
                        SELECT h.id
                        FROM chat_handle_join chj
                        JOIN handle h ON chj.handle_id = h.ROWID
                        WHERE chj.chat_id = (
                            SELECT chat_id FROM chat_message_join WHERE message_id = ?
                        )
                        LIMIT 1
                    """
                    cursor.execute(temp_query, (rowid,))
                    result = cursor.fetchone()
                    participant = result[0] if result else f"Unknown_{rowid}"
                else:
                    participant = handle_identifier

                # Try to resolve to contact name if available
                if contact_lookup and participant and not participant.startswith("Unknown_"):
                    contact_name = contact_lookup.get_contact_name(participant, show_phone=False)
                    if contact_name:
                        participant = contact_name

            # Initialize participant stats if needed
            if participant not in participant_stats:
                participant_stats[participant] = {
                    'message_count': 0,
                    'sent_by_me': 0,
                    'received_from_participant': 0,
                    'total_response_time': 0,
                    'response_count': 0,
                    'avg_response_time': 0,
                    'median_response_time': 0,
                    'response_times': []
                }

            # Track sent vs received counts
            # For forensic accuracy, always track by actual person, never by group chat name
            if is_from_me:
                # This is "Me" sending - track for the recipient
                # In 1-on-1 chats, the handle_identifier should always be set, but if it's NULL,
                # we need to look up the chat's participant handle from the chat itself
                if handle_identifier:
                    recipient = handle_identifier
                    if contact_lookup:
                        contact_name = contact_lookup.get_contact_name(handle_identifier, show_phone=False)
                        if contact_name:
                            recipient = contact_name
                else:
                    # Handle NULL handle_id for messages from Me - look up chat participant
                    # Since we filtered to 1-on-1 chats only, there should be exactly one participant handle
                    temp_query = """
                        SELECT h.id
                        FROM chat_handle_join chj
                        JOIN handle h ON chj.handle_id = h.ROWID
                        WHERE chj.chat_id = (
                            SELECT chat_id FROM chat_message_join WHERE message_id = ?
                        )
                        LIMIT 1
                    """
                    cursor.execute(temp_query, (rowid,))
                    result = cursor.fetchone()
                    recipient = result[0] if result else None
                    if recipient and contact_lookup:
                        contact_name = contact_lookup.get_contact_name(recipient, show_phone=False)
                        if contact_name:
                            recipient = contact_name

                if recipient:
                    # Initialize recipient stats if needed
                    if recipient not in participant_stats:
                        participant_stats[recipient] = {
                            'message_count': 0,
                            'sent_by_me': 0,
                            'received_from_participant': 0,
                            'total_response_time': 0,
                            'response_count': 0,
                            'avg_response_time': 0,
                            'median_response_time': 0,
                            'response_times': []
                        }

                    # Count this as a message sent TO this recipient
                    participant_stats[recipient]['sent_by_me'] += 1
                    participant_stats[recipient]['message_count'] += 1
            else:
                # Message received from this participant
                participant_stats[participant]['message_count'] += 1
                participant_stats[participant]['received_from_participant'] += 1

            # Calculate response time if this is a response to a different participant
            if prev_date and prev_is_from_me is not None:
                # Check if this is a response (different sender than previous)
                prev_participant = "Me" if prev_is_from_me else (prev_handle if prev_handle else "Unknown")

                if participant != prev_participant:
                    # Calculate response time in minutes
                    current_time = apple_timestamp_to_datetime(message_date or date)
                    previous_time = apple_timestamp_to_datetime(prev_date)

                    if current_time and previous_time:
                        # Apply timezone conversion
                        if tz_from:
                            try:
                                tz_from_zone = ZoneInfo(tz_from)
                                if current_time.tzinfo is None:
                                    current_time = current_time.replace(tzinfo=tz_from_zone)
                                    previous_time = previous_time.replace(tzinfo=tz_from_zone)
                            except Exception:
                                pass

                        if tz_to:
                            try:
                                tz_to_zone = ZoneInfo(tz_to)
                                current_time = current_time.astimezone(tz_to_zone)
                                previous_time = previous_time.astimezone(tz_to_zone)
                            except Exception:
                                pass

                        response_time_seconds = (current_time - previous_time).total_seconds()
                        response_time_minutes = response_time_seconds / 60

                        # Only count reasonable response times (within 24 hours)
                        if 0 < response_time_minutes < 1440:  # 24 hours in minutes
                            participant_stats[participant]['total_response_time'] += response_time_minutes
                            participant_stats[participant]['response_count'] += 1
                            participant_stats[participant]['response_times'].append(response_time_minutes)
                            response_times.append(response_time_minutes)

        # Calculate averages and percentages
        for participant in participant_stats:
            stats = participant_stats[participant]
            stats['percentage'] = (stats['message_count'] / total_messages_1on1) * 100 if total_messages_1on1 > 0 else 0
            if stats['response_count'] > 0:
                stats['avg_response_time'] = stats['total_response_time'] / stats['response_count']
                stats['median_response_time'] = statistics.median(stats['response_times'])
            # Remove the raw response_times list before returning (we don't need it in the output)
            del stats['response_times']

        # Sort by message count (most active first)
        sorted_participants = sorted(
            participant_stats.items(),
            key=lambda x: x[1]['message_count'],
            reverse=True
        )

        # Calculate overall response time stats
        overall_avg_response = sum(response_times) / len(response_times) if response_times else 0
        overall_median_response = statistics.median(response_times) if response_times else 0

        # Check for orphaned messages (messages not linked to any chat) - forensically important
        orphaned_query = """
            SELECT COUNT(DISTINCT m.ROWID)
            FROM message m
            LEFT JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
            WHERE cmj.message_id IS NULL
        """
        cursor.execute(orphaned_query)
        orphaned_count = cursor.fetchone()[0]

        activity_info = {
            'total_messages': total_messages_all,  # Total of ALL messages (including group chats)
            'total_messages_1on1': total_messages_1on1,  # Messages in 1-on-1 chats only
            'orphaned_messages': orphaned_count,  # Messages not linked to any chat
            'participant_stats': dict(sorted_participants),
            'overall_avg_response_minutes': overall_avg_response,
            'overall_median_response_minutes': overall_median_response,
            'timezone_note': f" (times in {tz_to})" if tz_to else " (times in UTC)"
        }

        conn.close()
        return activity_info

    except Exception as exc:
        print(f"Error analyzing participant activity: {exc}")
        if 'conn' in locals():
            conn.close()
        return None

# --- Wrapper: turn activity rows into pair counts for the graph -------------
def get_participant_activity_pairs(
    db_path,
    contact_lookup=None,
    tz_from="UTC",
    tz_to=None,
    chat_id=None,   # optional: pass to scope to a single chat if you like
):
    """
    Aggregate person<->person pairs from raw message data.
    Returns a list of dicts: {'sender': A, 'recipient': B, 'message_count': N}
    """
    try:
        import sqlite3
        from collections import defaultdict

        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Build query to get raw messages with previous sender info
        if chat_id:
            query = """
                SELECT
                    m.ROWID,
                    m.is_from_me,
                    h.id as handle_identifier,
                    cmj.chat_id,
                    cmj.message_date,
                    LAG(m.is_from_me) OVER (PARTITION BY cmj.chat_id ORDER BY cmj.message_date) as prev_is_from_me,
                    LAG(h.id) OVER (PARTITION BY cmj.chat_id ORDER BY cmj.message_date) as prev_handle
                FROM message m
                JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
                LEFT JOIN handle h ON m.handle_id = h.ROWID
                WHERE cmj.chat_id = ?
                ORDER BY cmj.chat_id, cmj.message_date ASC
            """
            cursor.execute(query, (chat_id,))
        else:
            query = """
                SELECT
                    m.ROWID,
                    m.is_from_me,
                    h.id as handle_identifier,
                    cmj.chat_id,
                    cmj.message_date,
                    LAG(m.is_from_me) OVER (PARTITION BY cmj.chat_id ORDER BY cmj.message_date) as prev_is_from_me,
                    LAG(h.id) OVER (PARTITION BY cmj.chat_id ORDER BY cmj.message_date) as prev_handle
                FROM message m
                JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
                LEFT JOIN handle h ON m.handle_id = h.ROWID
                ORDER BY cmj.chat_id, cmj.message_date ASC
            """
            cursor.execute(query)

        messages = cursor.fetchall()
        conn.close()

        if not messages:
            return []

        # 2) helpers to normalize names consistently with the rest of your app
        def norm_person(is_from_me, handle_identifier):
            if is_from_me:
                return "Me"

            # Try to resolve to contact name if available
            if contact_lookup and handle_identifier:
                contact_name = contact_lookup.get_contact_name(handle_identifier, show_phone=False)
                if contact_name:
                    return contact_name

            return handle_identifier or "Unknown"

        # 3) accumulate undirected pair counts
        pair_counts = defaultdict(int)

        # Process messages
        for msg in messages:
            rowid, is_from_me, handle_identifier, chat_id_val, message_date, prev_is_from_me, prev_handle = msg

            # Get current sender
            sender = norm_person(is_from_me, handle_identifier)

            # Get previous sender if available
            if prev_is_from_me is not None:
                prev_sender = norm_person(prev_is_from_me, prev_handle)

                # Only count an interaction when the speaker switches (A then B)
                if prev_sender and sender and prev_sender != sender:
                    # undirected pair (A,B) == (B,A)
                    a, b = sorted([prev_sender, sender], key=lambda s: s.lower())
                    pair_counts[(a, b)] += 1

        # 4) convert to the graph builder's simple shape
        result = [
            {"sender": a, "recipient": b, "message_count": cnt}
            for (a, b), cnt in sorted(pair_counts.items(), key=lambda kv: (-kv[1], kv[0][0].lower(), kv[0][1].lower()))
        ]
        return result

    except Exception as exc:
        print(f"Error in get_participant_activity_pairs: {exc}")
        return []


def get_participant_interaction_matrix(db_path, contact_lookup=None, tz_from="UTC", tz_to=None):
    """Analyze participant interaction patterns - who talks to whom most often"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Get all chats with their participants (without message counts - we'll count those accurately later)
        query = """
            SELECT
                c.ROWID as chat_id,
                c.display_name,
                c.chat_identifier,
                GROUP_CONCAT(DISTINCT h.id) as participant_handles
            FROM chat c
            LEFT JOIN chat_handle_join chj ON c.ROWID = chj.chat_id
            LEFT JOIN handle h ON chj.handle_id = h.ROWID
            GROUP BY c.ROWID, c.display_name, c.chat_identifier
        """
        cursor.execute(query)
        chat_data = cursor.fetchall()

        # Initialize interaction matrix
        interaction_matrix = {}
        total_interactions = 0
        chat_details = []

        # Process each chat to build interaction patterns
        for chat_info in chat_data:
            chat_id, display_name, chat_identifier, participant_handles_str = chat_info

            # Get detailed message data for this chat - count messages accurately
            msg_query = """
                SELECT
                    m.is_from_me,
                    h.id as handle_identifier,
                    m.date
                FROM message m
                JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
                LEFT JOIN handle h ON m.handle_id = h.ROWID
                WHERE cmj.chat_id = ?
                ORDER BY cmj.message_date ASC
            """
            cursor.execute(msg_query, (chat_id,))
            messages = cursor.fetchall()

            # Skip chats with no messages
            if not messages:
                continue

            # Get actual message count
            message_count = len(messages)

            if not participant_handles_str:
                continue

            participant_handles = participant_handles_str.split(',')

            # Resolve participant names
            participants = []
            for handle in participant_handles:
                if handle:
                    if contact_lookup:
                        contact_name = contact_lookup.get_contact_name(handle, show_phone=False)
                        participant_name = contact_name if contact_name else handle
                    else:
                        participant_name = handle
                    participants.append(participant_name)

            # Add "Me" since all chats involve the phone owner
            participants.append("Me")
            participants = sorted(set(participants))

            # Create chat details for display
            chat_detail = {
                'chat_id': chat_id,
                'display_name': display_name if display_name else f"Chat {chat_id}",
                'participants': participants,
                'message_count': message_count,
                'participant_count': len(participants)
            }
            chat_details.append(chat_detail)

            # Count interactions in this chat
            participant_message_counts = {}
            for msg in messages:
                is_from_me, handle_identifier, date = msg

                if is_from_me:
                    sender = "Me"
                else:
                    if contact_lookup and handle_identifier:
                        contact_name = contact_lookup.get_contact_name(handle_identifier, show_phone=False)
                        sender = contact_name if contact_name else handle_identifier
                    else:
                        sender = handle_identifier if handle_identifier else "Unknown"

                if sender not in participant_message_counts:
                    participant_message_counts[sender] = 0
                participant_message_counts[sender] += 1

            # Build interaction pairs for this chat
            # Only count direct 1-on-1 interactions for forensic accuracy
            # For 1-on-1 chats: there are exactly 2 participants (Me and one other)
            # For group chats: skip them to avoid inflated/misleading counts

            if len(participants) == 2:
                # This is a 1-on-1 chat - count all messages toward this interaction pair
                pair = tuple(sorted(participants))
                if pair not in interaction_matrix:
                    interaction_matrix[pair] = {
                        'message_count': 0,
                        'chats': []
                    }
                interaction_matrix[pair]['message_count'] += message_count
                if chat_id not in [c['chat_id'] for c in interaction_matrix[pair]['chats']]:
                    interaction_matrix[pair]['chats'].append({
                        'chat_id': chat_id,
                        'display_name': chat_detail['display_name'],
                        'messages_in_chat': message_count
                    })

        conn.close()

        if not interaction_matrix:
            return None

        # Sort interactions by message count
        sorted_interactions = sorted(
            interaction_matrix.items(),
            key=lambda x: x[1]['message_count'],
            reverse=True
        )

        # Calculate total for percentages
        total_interactions = sum(data['message_count'] for _, data in sorted_interactions)

        # Format results
        interaction_info = {
            'total_interactions': total_interactions,
            'interaction_pairs': [],
            'chat_summary': sorted(chat_details, key=lambda x: x['message_count'], reverse=True)[:10]  # Top 10 chats
        }

        # Build formatted interaction pairs
        for (person1, person2), data in sorted_interactions[:20]:  # Top 20 interactions
            percentage = (data['message_count'] / total_interactions * 100) if total_interactions > 0 else 0
            interaction_info['interaction_pairs'].append({
                'person1': person1,
                'person2': person2,
                'message_count': data['message_count'],
                'percentage': percentage,
                'chat_count': len(data['chats']),
                'chats': data['chats'][:5]  # Show top 5 chats for this pair
            })

        return interaction_info

    except Exception as exc:
        print(f"Error analyzing participant interactions: {exc}")
        return None

def get_communication_heatmap_info(db_path, chat_id=None, tz_from="UTC", tz_to=None, date_from_apple=None, date_to_apple=None):
    """Analyze communication patterns by time of day and day of week"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Build query conditions for filtering
        query_conditions = []
        query_params = []

        if chat_id:
            query_conditions.append("cmj.chat_id = ?")
            query_params.append(chat_id)

        if date_from_apple:
            query_conditions.append("cmj.message_date >= ?")
            query_params.append(date_from_apple)

        if date_to_apple:
            query_conditions.append("cmj.message_date <= ?")
            query_params.append(date_to_apple)

        where_clause = ""
        if query_conditions:
            where_clause = "WHERE " + " AND ".join(query_conditions)

        # Build query for message timestamps with date filtering
        query = f"""
            SELECT
                m.date,
                m.is_from_me,
                h.id as handle_identifier
            FROM message m
            JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
            LEFT JOIN handle h ON m.handle_id = h.ROWID
            {where_clause}
            ORDER BY m.date ASC
        """
        cursor.execute(query, query_params)

        messages = cursor.fetchall()
        conn.close()

        if not messages:
            return None

        # Initialize timezone conversion
        tz_from_zone = None
        tz_to_zone = None

        if tz_from:
            try:
                from zoneinfo import ZoneInfo
                tz_from_zone = ZoneInfo(tz_from)
            except Exception:
                tz_from_zone = None

        if tz_to:
            try:
                from zoneinfo import ZoneInfo
                tz_to_zone = ZoneInfo(tz_to)
            except Exception:
                tz_to_zone = tz_from_zone

        # Initialize heatmap data structures
        # Hour of day (0-23)
        hourly_activity = [0] * 24
        # Day of week (0=Monday, 6=Sunday)
        daily_activity = [0] * 7
        # Combined heatmap: [day_of_week][hour]
        combined_heatmap = [[0 for _ in range(24)] for _ in range(7)]

        day_names = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
        total_messages = 0

        # Process each message
        for msg in messages:
            date, is_from_me, handle_identifier = msg

            # Convert Apple timestamp to datetime
            dt = apple_timestamp_to_datetime(date) if date else None
            if not dt:
                continue

            # Apply timezone conversion
            if tz_from_zone and dt.tzinfo is None:
                dt = dt.replace(tzinfo=tz_from_zone)
            elif dt.tzinfo is None:
                # Default to UTC if no timezone info
                from zoneinfo import ZoneInfo
                dt = dt.replace(tzinfo=ZoneInfo('UTC'))

            if tz_to_zone and dt.tzinfo:
                dt = dt.astimezone(tz_to_zone)

            # Extract time patterns
            hour = dt.hour
            day_of_week = dt.weekday()  # 0=Monday, 6=Sunday

            # Increment counters
            hourly_activity[hour] += 1
            daily_activity[day_of_week] += 1
            combined_heatmap[day_of_week][hour] += 1
            total_messages += 1

        # Find peak activity periods
        peak_hour = hourly_activity.index(max(hourly_activity))
        peak_day = daily_activity.index(max(daily_activity))

        # Calculate percentages
        hourly_percentages = [(count / total_messages * 100) if total_messages > 0 else 0 for count in hourly_activity]
        daily_percentages = [(count / total_messages * 100) if total_messages > 0 else 0 for count in daily_activity]

        # Find quietest periods
        quietest_hour = hourly_activity.index(min(hourly_activity))
        quietest_day = daily_activity.index(min(daily_activity))

        # Generate hourly summary (group into 6 4-hour blocks)
        time_blocks = [
            ("Night", 0, 4),      # 12 AM - 4 AM
            ("Early Morning", 4, 8),  # 4 AM - 8 AM
            ("Morning", 8, 12),   # 8 AM - 12 PM
            ("Afternoon", 12, 16), # 12 PM - 4 PM
            ("Evening", 16, 20),  # 4 PM - 8 PM
            ("Late Evening", 20, 24)  # 8 PM - 12 AM
        ]

        block_activity = []
        for block_name, start_hour, end_hour in time_blocks:
            block_total = sum(hourly_activity[start_hour:end_hour])
            block_percentage = (block_total / total_messages * 100) if total_messages > 0 else 0
            block_activity.append({
                'name': block_name,
                'hours': f"{start_hour:02d}:00-{end_hour:02d}:00",
                'count': block_total,
                'percentage': block_percentage
            })

        # Sort time blocks by activity
        block_activity.sort(key=lambda x: x['count'], reverse=True)

        heatmap_info = {
            'total_messages': total_messages,
            'timezone_used': tz_to if tz_to else tz_from if tz_from else "UTC",
            'peak_hour': peak_hour,
            'peak_hour_count': hourly_activity[peak_hour],
            'peak_day': day_names[peak_day],
            'peak_day_count': daily_activity[peak_day],
            'quietest_hour': quietest_hour,
            'quietest_hour_count': hourly_activity[quietest_hour],
            'quietest_day': day_names[quietest_day],
            'quietest_day_count': daily_activity[quietest_day],
            'hourly_activity': list(zip(range(24), hourly_activity, hourly_percentages)),
            'daily_activity': list(zip(day_names, daily_activity, daily_percentages)),
            'time_blocks': block_activity,
            'combined_heatmap': combined_heatmap,
            'day_names': day_names
        }

        return heatmap_info

    except Exception as exc:
        print(f"Error analyzing communication heatmap: {exc}")
        return None

def get_deleted_message_info(db_path, chat_id=None, tz_from="UTC", tz_to=None):
    """Analyze deleted messages by detecting ROWID gaps and estimating deletion times"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Build query to analyze ROWID gaps, optionally filtered by chat
        if chat_id:
            query = """
                SELECT
                    m.ROWID,
                    LAG(m.ROWID) OVER (ORDER BY m.ROWID) as prev_rowid,
                    m.date,
                    LAG(m.date) OVER (ORDER BY m.ROWID) as prev_date,
                    cmj.chat_id
                FROM message m
                JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
                WHERE cmj.chat_id = ?
                ORDER BY m.ROWID
            """
            cursor.execute(query, (chat_id,))
        else:
            query = """
                SELECT
                    m.ROWID,
                    LAG(m.ROWID) OVER (ORDER BY m.ROWID) as prev_rowid,
                    m.date,
                    LAG(m.date) OVER (ORDER BY m.ROWID) as prev_date,
                    NULL as chat_id
                FROM message m
                ORDER BY m.ROWID
            """
            cursor.execute(query)

        # Analyze gaps
        gaps = []
        total_deleted = 0

        for row in cursor.fetchall():
            rowid, prev_rowid, date, prev_date, chat_filter = row
            if prev_rowid and rowid != prev_rowid + 1:
                gap_size = rowid - prev_rowid - 1
                total_deleted += gap_size

                # Convert timestamps for time estimation
                prev_dt = apple_timestamp_to_datetime(prev_date) if prev_date else None
                curr_dt = apple_timestamp_to_datetime(date) if date else None

                # Apply timezone conversion if specified
                if prev_dt and curr_dt:
                    if tz_from:
                        try:
                            tz_from_zone = ZoneInfo(tz_from)
                            if prev_dt.tzinfo is None:
                                prev_dt = prev_dt.replace(tzinfo=tz_from_zone)
                                curr_dt = curr_dt.replace(tzinfo=tz_from_zone)
                        except Exception:
                            pass

                    if tz_to:
                        try:
                            tz_to_zone = ZoneInfo(tz_to)
                            prev_dt = prev_dt.astimezone(tz_to_zone)
                            curr_dt = curr_dt.astimezone(tz_to_zone)
                        except Exception:
                            pass

                gaps.append({
                    'prev_rowid': prev_rowid,
                    'next_rowid': rowid,
                    'gap_size': gap_size,
                    'prev_time': prev_dt.strftime("%Y-%m-%d %H:%M:%S") if prev_dt else 'Unknown',
                    'next_time': curr_dt.strftime("%Y-%m-%d %H:%M:%S") if curr_dt else 'Unknown'
                })

        conn.close()

        # Sort gaps by size (largest first) for display
        gaps.sort(key=lambda x: x['gap_size'], reverse=True)

        deleted_info = {
            'total_deleted': total_deleted,
            'gap_count': len(gaps),
            'largest_gaps': gaps[:5],  # Top 5 largest gaps
            'timezone_note': f" (times in {tz_to})" if tz_to else " (times in UTC)"
        }

        return deleted_info

    except Exception as exc:
        print(f"Error analyzing deleted messages: {exc}")
        return None

def get_epoch_info(db_path, chat_id=None):
    """Get epoch timestamp information for the report"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Get sample timestamps
        if chat_id:
            query = """
                SELECT MIN(cmj.message_date), MAX(cmj.message_date)
                FROM chat_message_join cmj
                WHERE cmj.chat_id = ? AND cmj.message_date IS NOT NULL
            """
            cursor.execute(query, (chat_id,))
        else:
            query = """
                SELECT MIN(cmj.message_date), MAX(cmj.message_date)
                FROM chat_message_join cmj
                WHERE cmj.message_date IS NOT NULL
            """
            cursor.execute(query)

        result = cursor.fetchone()
        conn.close()

        if result and result[0]:
            sample_timestamp = result[0]

            # Create epoch info
            epoch_info = {
                'format': 'Apple/Cocoa Timestamp',
                'epoch_start': 'January 1, 2001 00:00:00 UTC',
                'resolution': 'Nanoseconds',
                'sample_raw': sample_timestamp,
                'sample_converted': apple_timestamp_to_datetime(sample_timestamp).strftime("%Y-%m-%d %H:%M:%S.%f UTC") if apple_timestamp_to_datetime(sample_timestamp) else 'N/A'
            }
            return epoch_info

        return None

    except Exception as exc:
        print(f"Error getting epoch info: {exc}")
        return None

def get_message_date_range(db_path, chat_id=None, tz_from="UTC", tz_to=None):
    """Get the date range (first and last message timestamps) from the database with timezone conversion"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Build query to get min/max message dates, optionally filtered by chat_id
        if chat_id:
            query = """
                SELECT MIN(cmj.message_date), MAX(cmj.message_date)
                FROM chat_message_join cmj
                WHERE cmj.chat_id = ? AND cmj.message_date IS NOT NULL
            """
            cursor.execute(query, (chat_id,))
        else:
            query = """
                SELECT MIN(cmj.message_date), MAX(cmj.message_date)
                FROM chat_message_join cmj
                WHERE cmj.message_date IS NOT NULL
            """
            cursor.execute(query)

        result = cursor.fetchone()
        conn.close()

        if result and result[0] and result[1]:
            min_timestamp, max_timestamp = result

            # Convert Apple timestamps to datetime objects
            min_date = apple_timestamp_to_datetime(min_timestamp)
            max_date = apple_timestamp_to_datetime(max_timestamp)

            if min_date and max_date:
                # Apply timezone conversion if specified
                if tz_from:
                    try:
                        tz_from_zone = ZoneInfo(tz_from)
                        if min_date.tzinfo is None:
                            min_date = min_date.replace(tzinfo=tz_from_zone)
                            max_date = max_date.replace(tzinfo=tz_from_zone)
                    except Exception:
                        pass  # Use original timezone if conversion fails

                if tz_to:
                    try:
                        tz_to_zone = ZoneInfo(tz_to)
                        min_date = min_date.astimezone(tz_to_zone)
                        max_date = max_date.astimezone(tz_to_zone)
                    except Exception:
                        pass  # Use original timezone if conversion fails

                # Return as YYYY-MM-DD format strings
                return min_date.strftime("%Y-%m-%d"), max_date.strftime("%Y-%m-%d")

        return None, None

    except Exception as exc:
        print(f"Error getting message date range: {exc}")
        return None, None

def collect_chat_ids(db_path, contact_lookup: ContactLookup = None):
    """Get list of available chat IDs from the database with contact name resolution"""
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        # Get chat info with message counts and participant info
        cursor.execute("""
            SELECT c.ROWID, c.chat_identifier, c.display_name, COUNT(cmj.message_id) as msg_count
            FROM chat c
            LEFT JOIN chat_message_join cmj ON c.ROWID = cmj.chat_id
            GROUP BY c.ROWID, c.chat_identifier, c.display_name
            HAVING msg_count > 0
            ORDER BY msg_count DESC
        """)

        chats = cursor.fetchall()

        chat_list = []
        for chat_id, identifier, display_name, msg_count in chats:
            name = display_name if display_name else identifier

            # Try to enhance with contact names if available
            enhanced_name = name
            if contact_lookup and identifier:
                # For single-person chats, try to resolve the phone number
                if not display_name and identifier:
                    contact_name = contact_lookup.get_contact_name(identifier, show_phone=False)
                    if contact_name:
                        enhanced_name = f"{contact_name} ({identifier})"
                    else:
                        enhanced_name = identifier

            chat_list.append(f"{chat_id}: {enhanced_name} ({msg_count} messages)")

        conn.close()
        return chat_list, None

    except Exception as exc:
        return [], str(exc)


def read_messages_from_db(
    db_path,
    contact_lookup: ContactLookup,
    *,
    chat_id=None,
    tz_from="UTC",
    tz_to=None,
    show_tz=False,
    show_phone=False,
    date_from=None,
    date_to=None,
    search_text=None,
):
    """Read messages from SQLite database and return structured data with contact resolution"""

    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # Parse date filters
    tz_from_zone = None
    if tz_from:
        try:
            tz_from_zone = ZoneInfo(tz_from)
        except Exception:
            tz_from_zone = None

    # For date filtering, use destination timezone if available, otherwise source timezone
    # This is because users enter dates in their preferred display timezone (destination)
    date_filter_zone = None
    if tz_to:
        try:
            date_filter_zone = ZoneInfo(tz_to)
        except Exception:
            date_filter_zone = tz_from_zone
    else:
        date_filter_zone = tz_from_zone

    date_from_apple = None
    date_to_apple = None
    if date_from:
        parsed_date = parse_date_filter(date_from, date_filter_zone)
        if parsed_date:
            date_from_apple = datetime_to_apple_timestamp(parsed_date)

    if date_to:
        parsed_date = parse_date_filter(date_to, date_filter_zone)
        if parsed_date:
            # For end date, set to end of day if no time specified
            if ' ' not in date_to:
                parsed_date = parsed_date.replace(hour=23, minute=59, second=59)
            date_to_apple = datetime_to_apple_timestamp(parsed_date)

    # Build query with date filtering
    query_conditions = []
    query_params = []

    if chat_id:
        query_conditions.append("cmj.chat_id = ?")
        query_params.append(chat_id)

    if date_from_apple:
        query_conditions.append("COALESCE(cmj.message_date, m.date) >= ?")
        query_params.append(date_from_apple)

    if date_to_apple:
        query_conditions.append("COALESCE(cmj.message_date, m.date) <= ?")
        query_params.append(date_to_apple)

    if search_text:
        query_conditions.append("m.text LIKE ? COLLATE NOCASE")
        query_params.append(f"%{search_text}%")

    where_clause = ""
    if query_conditions:
        where_clause = "WHERE " + " AND ".join(query_conditions)

    print(f"\n🔍 QUERY DEBUG:")
    print(f"   Query conditions: {query_conditions}")
    print(f"   Query params: {query_params}")
    print(f"   Where clause: {where_clause if where_clause else 'None'}")

    # Modified query to include orphaned messages (messages not linked in chat_message_join)
    # These are forensically important as they exist in the database
    query = f"""
        SELECT
            m.ROWID,
            m.text,
            m.date,
            m.is_from_me,
            m.handle_id,
            h.id as handle_identifier,
            c.chat_identifier,
            c.display_name,
            COALESCE(cmj.message_date, m.date) as message_date,
            GROUP_CONCAT(a.filename, '|||') as filenames,
            GROUP_CONCAT(a.mime_type, '|||') as mime_types,
            GROUP_CONCAT(a.transfer_name, '|||') as transfer_names,
            MAX(CASE WHEN a.is_sticker = 1 THEN 1 ELSE 0 END) as is_sticker
        FROM message m
        LEFT JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
        LEFT JOIN chat c ON cmj.chat_id = c.ROWID
        LEFT JOIN handle h ON m.handle_id = h.ROWID
        LEFT JOIN message_attachment_join maj ON m.ROWID = maj.message_id
        LEFT JOIN attachment a ON maj.attachment_id = a.ROWID
        {where_clause}
        GROUP BY m.ROWID, m.text, m.date, m.is_from_me, m.handle_id, h.id, c.chat_identifier, c.display_name, COALESCE(cmj.message_date, m.date)
        ORDER BY COALESCE(cmj.message_date, m.date) ASC
    """

    cursor.execute(query, query_params)

    messages = cursor.fetchall()
    total_messages = len(messages)

    # Debug: Check for orphaned messages (forensically important)
    cursor.execute("SELECT COUNT(*) FROM message")
    db_total_messages = cursor.fetchone()[0]

    print(f"\n📊 MESSAGE COUNT SUMMARY:")
    print(f"   Total in database: {db_total_messages}")
    print(f"   Returned by query: {total_messages}")

    # Check for orphaned messages
    cursor.execute("""
        SELECT COUNT(*)
        FROM message m
        LEFT JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
        WHERE cmj.message_id IS NULL
    """)
    orphaned_count = cursor.fetchone()[0]

    if orphaned_count > 0:
        print(f"\nℹ️  ORPHANED MESSAGES:")
        print(f"   Found {orphaned_count} message(s) not linked to any chat")
        if total_messages == db_total_messages:
            print(f"   ✓ Orphaned messages ARE included in the query results")
        else:
            print(f"   ⚠️  WARNING: Orphaned messages NOT included! Missing {db_total_messages - total_messages} message(s)")

        # Get details of orphaned messages
        cursor.execute("""
            SELECT m.ROWID, m.text, m.date, m.is_from_me
            FROM message m
            LEFT JOIN chat_message_join cmj ON m.ROWID = cmj.message_id
            WHERE cmj.message_id IS NULL
            LIMIT 5
        """)
        orphaned_msgs = cursor.fetchall()
        print(f"   Sample orphaned messages:")
        for msg in orphaned_msgs:
            rowid, text, date, is_from_me = msg
            preview = (text[:50] + '...') if text and len(text) > 50 else text
            print(f"     - Message {rowid}: {'Sent' if is_from_me else 'Received'}, text='{preview}'")

    # Setup timezone conversion (tz_from_zone already set above for date filtering)
    tz_to_zone = None
    if tz_to:
        try:
            tz_to_zone = ZoneInfo(tz_to)
        except Exception:
            tz_to_zone = None


    # Process messages
    rows = []
    participants = set()

    for msg in tqdm(messages, desc="Processing messages", unit="msg", total=total_messages):
        msg_id, text, date, is_from_me, handle_id, handle_identifier, chat_identifier, display_name, message_date, filenames, mime_types, transfer_names, is_sticker = msg

        # Determine direction
        direction = "sent" if is_from_me else "received"

        # Process text
        text = text if text else ""

        # Process timestamp - use message_date from join table as it's more reliable
        timestamp_to_use = message_date if message_date else date
        dt_parsed = apple_timestamp_to_datetime(timestamp_to_use)

        ts_display = ""
        if dt_parsed:
            if dt_parsed.tzinfo is None and tz_from_zone:
                dt_parsed = dt_parsed.replace(tzinfo=tz_from_zone)
            if tz_to_zone:
                dt_parsed = dt_parsed.astimezone(tz_to_zone)
            ts_display = dt_parsed.strftime("%m/%d/%Y %I:%M %p")

        if show_tz and tz_to:
            ts_display = f"{ts_display} {tz_to}"

        # Check if this is an orphaned message (not linked to any chat)
        is_orphaned = chat_identifier is None and display_name is None

        # Handle sender/recipient info with contact resolution
        sender = handle_identifier if handle_identifier else f"Handle {handle_id}"
        recipient = "Me"

        # Try to resolve contact name
        contact_name = None
        if handle_identifier:
            contact_name = contact_lookup.get_contact_name(handle_identifier, show_phone=show_phone)

        if direction == "sent":
            sender = "Me"
            recipient = handle_identifier if handle_identifier else f"Handle {handle_id}"

        # Use contact name for display if available
        sender_display = sender
        recipient_display = recipient

        if direction == "sent" and contact_name:
            recipient_display = contact_name
        elif direction == "received" and contact_name:
            sender_display = contact_name

        # Mark orphaned messages for forensic tracking (after contact resolution)
        if is_orphaned:
            orphan_tag = " [ORPHANED - No Chat Link]"
            if direction == "sent":
                recipient_display = recipient_display + orphan_tag
            else:
                sender_display = sender_display + orphan_tag


        # Track participants
        if handle_identifier:
            participants.add(contact_name if contact_name else handle_identifier)

        # Parse concatenated attachment data (use first attachment for backward compatibility)
        filename_list = filenames.split('|||') if filenames else [None]
        mime_type_list = mime_types.split('|||') if mime_types else [None]
        transfer_name_list = transfer_names.split('|||') if transfer_names else [None]

        # Use first attachment for the row (backward compatible with existing code)
        filename = filename_list[0] if filename_list[0] else None
        mime_type = mime_type_list[0] if mime_type_list[0] else None
        transfer_name = transfer_name_list[0] if transfer_name_list[0] else None

        rows.append({
            "direction": direction,
            "sender": sender,
            "recipient": recipient,
            "sender_display": sender_display,
            "recipient_display": recipient_display,
            "timestamp": ts_display,
            "text": text,
            "attachment_filename": filename,
            "attachment_mime_type": mime_type,
            "attachment_transfer_name": transfer_name,
            "attachment_is_sticker": is_sticker,
            "chat_display_name": display_name,
            "chat_identifier": chat_identifier,
        })

    conn.close()

    # Get chat info for metadata
    chat_info = {"identifier": "All chats", "display_name": "All conversations", "participant_count": len(participants)}
    if chat_id:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        cursor.execute("SELECT chat_identifier, display_name FROM chat WHERE ROWID = ?", (chat_id,))
        result = cursor.fetchone()
        if result:
            chat_info = {
                "identifier": result[0],
                "display_name": result[1] if result[1] else result[0],
                "participant_count": len(participants)
            }
        conn.close()

    return rows, participants, chat_info

def generate_pagination_nav(current_page, total_pages, page_filename_format="page_{}.html"):
    """Generate HTML navigation for pagination"""
    nav_html = '<div style="text-align: center; padding: 20px; background-color: #f5f5f5; border-radius: 10px; margin: 20px 0;">'
    nav_html += '<div style="font-size: 16px; margin-bottom: 10px;">'

    # Previous button
    if current_page > 1:
        prev_page = page_filename_format.format(current_page - 1)
        nav_html += f'<a href="{prev_page}" style="padding: 10px 20px; margin: 0 5px; background-color: #007AFF; color: white; text-decoration: none; border-radius: 5px;">← Previous</a>'
    else:
        nav_html += '<span style="padding: 10px 20px; margin: 0 5px; background-color: #ccc; color: #666; border-radius: 5px;">← Previous</span>'

    # Page indicator
    nav_html += f'<span style="padding: 10px 20px; margin: 0 10px; font-weight: bold;">Page {current_page} of {total_pages}</span>'

    # Next button
    if current_page < total_pages:
        next_page = page_filename_format.format(current_page + 1)
        nav_html += f'<a href="{next_page}" style="padding: 10px 20px; margin: 0 5px; background-color: #007AFF; color: white; text-decoration: none; border-radius: 5px;">Next →</a>'
    else:
        nav_html += '<span style="padding: 10px 20px; margin: 0 5px; background-color: #ccc; color: #666; border-radius: 5px;">Next →</span>'

    nav_html += '</div>'

    # Index link
    nav_html += '<div style="margin-top: 10px;">'
    nav_html += '<a href="index.html" style="color: #007AFF; text-decoration: none; font-size: 14px;">📋 Back to Overview</a>'
    nav_html += '</div>'

    nav_html += '</div>'
    return nav_html

def convert_imessage_db(
    db_path,
    out_path,
    *,
    addressbook_path=None,
    width=375,
    font_size=14,
    title="IMESSAGE MESSAGES",
    chat_id=None,
    contact_name=None,
    tz_from="UTC",
    tz_to=None,
    per_frame=0,
    show_tz=False,
    show_phone=False,
    powerpoint_path=None,
    date_from=None,
    date_to=None,
    show_epoch_info=False,
    show_deleted_info=False,
    show_participant_info=False,
    show_interaction_matrix=False,
    show_communication_heatmap=False,
    show_hash_verification=False,
    show_attachments=False,
    attachments_folder=None,
    search_text=None,
    messages_per_page=0,
    progress_callback: Optional[Callable[[str, int, int], None]] = None,
) -> Tuple[Path, List[Path]]:
    """Core conversion logic for SQLite database with contact resolution

    Args:
        progress_callback: Optional callback function(message: str, current: int, total: int)
                          Called to report progress. If None, no progress reporting.
    """
    db_path = Path(db_path)

    # Ensure out_path is valid
    if not out_path:
        raise ValueError("Output path cannot be empty")
    out_path = Path(out_path)

    # Helper function to safely call progress callback
    def report_progress(message: str, current: int = 0, total: int = 100):
        if progress_callback:
            try:
                progress_callback(message, current, total)
            except:
                pass  # Ignore callback errors

    # Hash verification - Calculate hashes before processing if requested
    pre_processing_hashes = None
    post_processing_hashes = None
    hash_verification_info = None

    # Hash verification - Calculate hashes before processing for both databases
    addressbook_pre_hashes = None
    addressbook_post_hashes = None

    if show_hash_verification:
        report_progress("Calculating pre-processing hashes...", 0, 100)
        print("Calculating pre-processing hashes for integrity verification...")
        pre_processing_hashes = calculate_file_hashes(db_path)
        if pre_processing_hashes:
            print(f"Pre-processing chat.db MD5: {pre_processing_hashes['md5']}")

        # Also hash AddressBook before processing if it exists
        if addressbook_path and Path(addressbook_path).exists():
            addressbook_pre_hashes = calculate_file_hashes(addressbook_path)
            if addressbook_pre_hashes:
                print(f"Pre-processing AddressBook MD5: {addressbook_pre_hashes['md5']}")

    # Initialize contact lookup
    report_progress("Initializing contact lookup...", 5, 100)
    contact_lookup = ContactLookup(addressbook_path)

    # Convert date filters to Apple timestamps for consistency across all queries
    date_from_apple = None
    date_to_apple = None
    if date_from or date_to:
        # Determine timezone for date filtering
        date_filter_zone = None
        if tz_to:
            try:
                date_filter_zone = ZoneInfo(tz_to)
            except Exception:
                pass
        if not date_filter_zone and tz_from:
            try:
                date_filter_zone = ZoneInfo(tz_from)
            except Exception:
                pass

        if date_from:
            parsed_date = parse_date_filter(date_from, date_filter_zone)
            if parsed_date:
                date_from_apple = datetime_to_apple_timestamp(parsed_date)

        if date_to:
            parsed_date = parse_date_filter(date_to, date_filter_zone)
            if parsed_date:
                # For end date, set to end of day if no time specified
                if ' ' not in date_to:
                    parsed_date = parsed_date.replace(hour=23, minute=59, second=59)
                date_to_apple = datetime_to_apple_timestamp(parsed_date)

    # Read messages from database
    report_progress("Reading messages from database...", 10, 100)
    rows, participants, chat_info = read_messages_from_db(
        db_path,
        contact_lookup,
        chat_id=chat_id,
        tz_from=tz_from,
        tz_to=tz_to,
        show_tz=show_tz,
        show_phone=show_phone,
        date_from=date_from,
        date_to=date_to,
        search_text=search_text,
    )
    report_progress(f"Loaded {len(rows)} messages", 30, 100)

    # Hash verification - Calculate hashes after processing if requested
    if show_hash_verification:
        print("Calculating post-processing hashes for integrity verification...")
        post_processing_hashes = calculate_file_hashes(db_path)
        if post_processing_hashes:
            print(f"Post-processing chat.db MD5: {post_processing_hashes['md5']}")

        # Also hash AddressBook after processing if it exists
        if addressbook_path and Path(addressbook_path).exists():
            addressbook_post_hashes = calculate_file_hashes(addressbook_path)
            if addressbook_post_hashes:
                print(f"Post-processing AddressBook MD5: {addressbook_post_hashes['md5']}")

        # Generate hash verification info
        hash_verification_info = get_hash_verification_info(
            db_path,
            addressbook_path,
            pre_processing_hashes,
            post_processing_hashes,
            addressbook_pre_hashes,
            addressbook_post_hashes
        )

        # Print integrity verification results
        if hash_verification_info:
            db_verified = hash_verification_info['database_hashes']['integrity_verified']
            ab_hashes = hash_verification_info.get('addressbook_hashes')
            ab_verified = ab_hashes['integrity_verified'] if ab_hashes else True

            if db_verified:
                print("CHAT.DB INTEGRITY VERIFIED - no changes detected during processing")
            else:
                print("WARNING: chat.db integrity check failed - database may have been modified")

            if ab_hashes:
                if ab_verified:
                    print("ADDRESSBOOK INTEGRITY VERIFIED - no changes detected during processing")
                else:
                    print("WARNING: AddressBook integrity check failed - database may have been modified")

            if db_verified and ab_verified:
                print("ALL DATABASES VERIFIED - forensic integrity maintained")
            elif not db_verified or not ab_verified:
                print("CAUTION: One or more databases failed integrity verification")

    # Format participants
    participants_display = ", ".join(sorted(p for p in participants if p)) if participants else "Unknown"

    # Get epoch info if requested
    epoch_rows = ""
    if show_epoch_info:
        epoch_info = get_epoch_info(db_path, chat_id)
        if epoch_info:
            epoch_rows = f"""
      <tr><th>Timestamp Format</th><td>{html.escape(epoch_info['format'])}</td></tr>
      <tr><th>Epoch Start</th><td>{html.escape(epoch_info['epoch_start'])}</td></tr>
      <tr><th>Resolution</th><td>{html.escape(epoch_info['resolution'])}</td></tr>
      <tr><th>Sample Raw Value</th><td>{epoch_info['sample_raw']}</td></tr>
      <tr><th>Sample Converted</th><td>{html.escape(epoch_info['sample_converted'])}</td></tr>"""

    # Get deleted message info if requested
    deleted_rows = ""
    if show_deleted_info:
        deleted_info = get_deleted_message_info(db_path, chat_id, tz_from, tz_to)
        if deleted_info:
            deleted_rows = f"""
      <tr><th>Total Deleted Messages</th><td>{deleted_info['total_deleted']}</td></tr>
      <tr><th>Number of Gaps</th><td>{deleted_info['gap_count']}</td></tr>"""

            # Add details about the largest gaps with individual timeframes
            if deleted_info['largest_gaps']:
                gap_details = []
                gap_timeframes = []

                for i, gap in enumerate(deleted_info['largest_gaps'], 1):
                    gap_detail = f"Gap {i}: {gap['gap_size']} messages between ROWID {gap['prev_rowid']} and {gap['next_rowid']}"
                    gap_details.append(gap_detail)

                    timeframe_detail = f"Gap {i}: {html.escape(gap['prev_time'])} to {html.escape(gap['next_time'])}"
                    gap_timeframes.append(timeframe_detail)

                # Join with <br> for HTML line breaks
                gaps_html = '<br>'.join(html.escape(detail) for detail in gap_details)
                timeframes_html = '<br>'.join(gap_timeframes)  # gap_timeframes already have escaped content

                deleted_rows += f"""
      <tr><th>Largest Gaps</th><td>{gaps_html}</td></tr>"""

                # Add individual time ranges for each gap
                deleted_rows += f"""
      <tr><th>Gap Timeframes</th><td>{timeframes_html}{html.escape(deleted_info['timezone_note'])}</td></tr>"""

    # Setup chat/thread info
    chat_row = ""
    chat_display_line = None
    title_suffix = ""

    if chat_id:
        chat_value = f"Chat {chat_id}"
        if contact_name:
            chat_value = f"{chat_value} ({contact_name})"
        elif chat_info["display_name"] and chat_info["display_name"] != chat_info["identifier"]:
            chat_value = f"{chat_value} ({chat_info['display_name']})"

        chat_row = f"<tr><th>Chat ID</th><td>{html.escape(chat_value)}</td></tr>"
        chat_display_line = f"Chat: {chat_value}"
        title_suffix = f"Chat {chat_id}"
    elif contact_name:
        chat_row = f"<tr><th>Contact</th><td>{html.escape(contact_name)}</td></tr>"
        chat_display_line = f"Contact: {contact_name}"
        title_suffix = contact_name

    # Add search info to title suffix
    if search_text:
        search_suffix = f"Search: '{search_text}'"
        if title_suffix:
            title_suffix = f"{title_suffix} - {search_suffix}"
        else:
            title_suffix = search_suffix

    phone_w = int(width)

    # Handle framing - separate logic for PowerPoint vs HTML
    report_progress("Preparing frames...", 35, 100)
    per_frame_val = max(0, int(per_frame)) if per_frame is not None else 0

    # First, create HTML frames using manual per_frame setting
    if per_frame_val:
        total_frames = math.ceil(len(rows) / per_frame_val) if rows else 1
        frames = [rows[i * per_frame_val : (i + 1) * per_frame_val] for i in range(total_frames)]
    else:
        frames = [rows]
        total_frames = 1

    # For PowerPoint, we'll calculate optimal frames separately later
    powerpoint_frames = None
    if powerpoint_path and rows:
        report_progress("Calculating PowerPoint frames...", 40, 100)
        # Calculate optimal frames for PowerPoint only
        optimal_per_frame, powerpoint_frames = calculate_optimal_messages_per_frame(rows, phone_w, font_size, chat_info=chat_info)
        print(f"Auto-calculated {optimal_per_frame} messages per frame for optimal PowerPoint fit ({len(powerpoint_frames)} slides)")
    title_display = title + (f" - {title_suffix}" if title_suffix else "")
    tz_line = f" (Times shown in {html.escape(str(tz_to))})" if tz_to else ""

    # Generate HTML - separate CSS and analysis sections
    # CSS-only header for individual pages
    css_head = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8"/>
<title>{html.escape(title_display)}</title>
<style>
  :root {{
    --card-bg: #ffffff;
    --card-shadow: 0 2px 8px rgba(0,0,0,0.1);
    --card-shadow-hover: 0 4px 16px rgba(0,0,0,0.15);
    --primary-color: #0b93f6;
    --secondary-color: #5856d6;
    --success-color: #34c759;
    --warning-color: #ff9500;
    --border-color: #e5e5e5;
    --text-primary: #1d1d1f;
    --text-secondary: #86868b;
    --bg-page: #f5f7fa;
    --header-gradient: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
  }}

  * {{ box-sizing: border-box; }}

  body {{
    margin: 0;
    padding: 0;
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
    background: var(--bg-page);
    color: var(--text-primary);
    line-height: 1.6;
  }}

  .main-container {{
    max-width: 1200px;
    margin: 0 auto;
    padding: 20px;
  }}

  .page-header {{
    background: var(--header-gradient);
    color: white;
    padding: 40px 20px;
    margin: -20px -20px 30px -20px;
    border-radius: 0 0 20px 20px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.15);
  }}

  .page-header h1 {{
    margin: 0 0 10px 0;
    font-size: 32px;
    font-weight: 700;
  }}

  .page-header .subtitle {{
    font-size: 16px;
    opacity: 0.95;
    margin: 5px 0;
  }}

  .metrics-grid {{
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
    gap: 20px;
    margin-bottom: 30px;
  }}

  .card {{
    background: var(--card-bg);
    border-radius: 12px;
    padding: 24px;
    box-shadow: var(--card-shadow);
    transition: transform 0.2s, box-shadow 0.2s;
    border: 1px solid var(--border-color);
  }}

  .card:hover {{
    transform: translateY(-2px);
    box-shadow: var(--card-shadow-hover);
  }}

  .card-header {{
    display: flex;
    align-items: center;
    margin-bottom: 20px;
    padding-bottom: 12px;
    border-bottom: 2px solid var(--border-color);
  }}

  .card-icon {{
    width: 40px;
    height: 40px;
    border-radius: 10px;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 20px;
    margin-right: 12px;
    font-weight: bold;
  }}

  .card-icon.primary {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; }}
  .card-icon.success {{ background: linear-gradient(135deg, #56ab2f 0%, #a8e063 100%); color: white; }}
  .card-icon.info {{ background: linear-gradient(135deg, #2193b0 0%, #6dd5ed 100%); color: white; }}
  .card-icon.warning {{ background: linear-gradient(135deg, #f09819 0%, #edde5d 100%); color: white; }}

  .card-title {{
    font-size: 18px;
    font-weight: 600;
    color: var(--text-primary);
    margin: 0;
  }}

  .card-subtitle {{
    font-size: 13px;
    color: var(--text-secondary);
    margin: 8px 0 0 0;
  }}

  .stat-row {{
    display: flex;
    justify-content: space-between;
    align-items: center;
    padding: 12px 0;
    border-bottom: 1px solid #f5f5f5;
  }}

  .stat-row:last-child {{ border-bottom: none; }}

  .stat-label {{
    font-size: 14px;
    color: var(--text-secondary);
    font-weight: 500;
  }}

  .stat-value {{
    font-size: 16px;
    font-weight: 600;
    color: var(--text-primary);
  }}

  .header-container {{
    width: 100%;
    margin-bottom: 24px;
  }}

  .phone-meta {{
    width: 100%;
    margin: 0;
  }}

  .chat-title {{
    font-weight: 700;
    font-size: 22px;
    margin-bottom: 8px;
    color: var(--text-primary);
  }}

  .subtle {{
    font-size: 14px;
    color: var(--text-secondary);
    margin: 4px 0 16px;
  }}

  table.meta {{
    width: 100%;
    border-collapse: collapse;
    font-size: 14px;
    margin: 0;
  }}

  table.meta th,
  table.meta td {{
    padding: 12px 16px;
    text-align: left;
    border-bottom: 1px solid var(--border-color);
  }}

  table.meta th {{
    background: #f9fafb;
    font-weight: 600;
    color: var(--text-primary);
    font-size: 13px;
    text-transform: uppercase;
    letter-spacing: 0.5px;
  }}

  table.meta tr:last-child td {{ border-bottom: none; }}

  .progress-bar-container {{
    width: 100%;
    height: 24px;
    background: #e9ecef;
    border-radius: 12px;
    overflow: hidden;
    position: relative;
    display: flex;
    align-items: center;
  }}

  .progress-bar-fill {{
    height: 100%;
    background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
    border-radius: 12px;
    transition: width 0.3s ease;
    display: flex;
    align-items: center;
    justify-content: flex-end;
    padding-right: 8px;
  }}

  .progress-bar-text {{
    position: absolute;
    right: 8px;
    font-size: 12px;
    font-weight: 600;
    color: var(--text-primary);
    z-index: 1;
  }}

  .progress-bar-fill-text {{
    font-size: 12px;
    font-weight: 600;
    color: white;
    white-space: nowrap;
  }}

  .messages-section {{
    background: var(--card-bg);
    border-radius: 12px;
    padding: 24px;
    box-shadow: var(--card-shadow);
    border: 1px solid var(--border-color);
    margin-top: 24px;
  }}

  .messages-section h2 {{
    font-size: 24px;
    font-weight: 700;
    margin: 0 0 20px 0;
    text-align: center;
    color: var(--text-primary);
  }}

  #phone-frame {{
    background: #fff;
    width: {phone_w}px;
    max-width: 100%;
    border: 6px solid #000;
    border-radius: 36px;
    overflow: hidden;
    box-shadow: 0 8px 24px rgba(0,0,0,0.2);
    margin: 0 auto 20px auto;
  }}
  #container {{ padding: 10px; }}
  .bubble {{
    border-radius: 18px;
    padding: 10px 14px;
    margin: 6px 8px;
    max-width: 75%;
    word-wrap: break-word;
    white-space: pre-wrap;
    font-size: {font_size}px;
    line-height: 1.35;
    position: relative;
  }}
  .bubble .meta-line {{
    font-size: {max(9, int(font_size * 0.79))}px;
    font-weight: 600;
    margin-bottom: 4px;
    opacity: 0.9;
  }}
  .bubble .ts {{
    display: block;
    font-size: {max(9, int(font_size * 0.79))}px;
    opacity: .75;
    margin-top: 6px;
  }}
  .sent {{ background: #0b93f6; color: #fff; margin-left: auto; }}
  .received {{ background: #e5e5ea; color: #000; margin-right: auto; }}
  .threaded-chat {{ display: flex; flex-direction: column; }}
  .frame-caption {{ text-align:center; font-size:12px; color:#555; margin-top:-6px; width: {phone_w}px; max-width: 100%; margin-left: auto; margin-right: auto; }}

  .page-break {{
    page-break-after: always;
    }}

</style>
</head>
<body>
<div class="main-container">
  <div class="page-header">
    <h1>{html.escape(title_display)}</h1>
    <div class="subtitle">📊 Total messages: {len(rows):,}{tz_line}</div>
    <div class="subtitle">👥 {html.escape(participants_display)}</div>
  </div>
"""

    # Full header with analysis sections for index.html and non-paginated output
    # Only show Conversation Details card if there's additional info beyond participants
    conversation_details_rows = []
    if chat_row:
        conversation_details_rows.append(chat_row)
    if per_frame_val:
        conversation_details_rows.append(f"<tr><th>Messages per frame</th><td>{per_frame_val}</td></tr>")
    if epoch_rows:
        conversation_details_rows.append(epoch_rows)

    head = css_head
    if conversation_details_rows:
        head += f"""
  <div class="card">
    <div class="card-header">
      <div class="card-icon primary">📋</div>
      <div>
        <div class="card-title">Conversation Details</div>
      </div>
    </div>
    <table class="meta phone-meta">
      {"".join(conversation_details_rows)}
    </table>
  </div>
"""

    # Add deleted messages as a separate card if requested
    if show_deleted_info and deleted_rows:
        head += f"""
  <div class="card">
    <div class="card-header">
      <div class="card-icon warning">⚠️</div>
      <div>
        <div class="card-title">Deleted Messages Analysis</div>
        <div class="card-subtitle">Analysis of message gaps and potential deletions</div>
      </div>
    </div>
    <table class="meta phone-meta">
{deleted_rows}
    </table>
  </div>
"""

    head += """
  <!-- Participant Activity Analysis Frame -->"""

    # Add participant activity analysis frame if requested
    if show_participant_info:
        activity_info = get_participant_activity_info(db_path, chat_id, contact_lookup, tz_from, tz_to, date_from_apple, date_to_apple)
        if activity_info:
            # Build orphaned message warning strings
            orphaned_count = activity_info.get('orphaned_messages', 0)
            orphaned_warning = f'<div style="font-size: 11px; color: #ff9500; margin-top: 4px;">⚠️ Includes {orphaned_count} orphaned message(s)</div>' if orphaned_count > 0 else ''
            orphaned_explanation = '<br>⚠️ <strong>Orphaned messages</strong> are messages that exist in the database but are not linked to any chat (forensically significant)' if orphaned_count > 0 else ''

            head += f"""
  <div class="card">
    <div class="card-header">
      <div class="card-icon success">📈</div>
      <div>
        <div class="card-title">Participant Activity Analysis</div>
        <div class="card-subtitle">Communication patterns and response times{activity_info['timezone_note']}</div>
      </div>
    </div>

    <div class="subtle" style="margin-bottom: 12px; font-style: italic;">
      ℹ️ Chart shows sent (by Me) vs received (from participant) message counts
    </div>

    <!-- Message Flow Chart -->
    <div style="margin-bottom: 30px;">
      <div style="font-weight: 600; margin-bottom: 15px; color: #333;">Messages Sent & Received by Participant</div>"""

            # Generate chart data for participants (excluding "Me" from the chart)
            chart_participants = [(p, s) for p, s in activity_info['participant_stats'].items() if p != "Me"]

            # Limit to top 10 for readability
            chart_participants = chart_participants[:20]

            for participant, stats in chart_participants:
                sent = stats['sent_by_me']
                received = stats['received_from_participant']
                total = sent + received

                sent_pct = (sent / total * 100) if total > 0 else 0
                received_pct = (received / total * 100) if total > 0 else 0

                head += f"""
      <div style="margin-bottom: 15px;">
        <div style="font-size: 13px; margin-bottom: 4px; color: #555;">
          <strong>{html.escape(participant)}</strong>
          <span style="color: #888;">({total:,} total: {sent:,} sent / {received:,} received)</span>
        </div>
        <div style="display: flex; height: 25px; background-color: #f0f0f0; border-radius: 4px; overflow: hidden;">
          <div style="width: {sent_pct:.1f}%; background-color: #007AFF; display: flex; align-items: center; justify-content: center; color: white; font-size: 11px; font-weight: 600;" title="Sent by Me: {sent:,}">
            {f'{sent_pct:.0f}%' if sent_pct > 10 else ''}
          </div>
          <div style="width: {received_pct:.1f}%; background-color: #34C759; display: flex; align-items: center; justify-content: center; color: white; font-size: 11px; font-weight: 600;" title="Received from {html.escape(participant)}: {received:,}">
            {f'{received_pct:.0f}%' if received_pct > 10 else ''}
          </div>
        </div>
      </div>"""

            head += f"""
      <div style="margin-top: 15px; padding: 10px; background-color: #f9f9f9; border-radius: 4px; font-size: 12px;">
        <div style="display: flex; gap: 20px; justify-content: center;">
          <div><span style="display: inline-block; width: 12px; height: 12px; background-color: #007AFF; border-radius: 2px; margin-right: 5px;"></span>Sent by Me</div>
          <div><span style="display: inline-block; width: 12px; height: 12px; background-color: #34C759; border-radius: 2px; margin-right: 5px;"></span>Received from Participant</div>
        </div>
      </div>
    </div>

    <div style="background-color: #f9f9f9; padding: 20px; border-radius: 8px; margin-bottom: 20px;">
      <div class="subtle" style="margin-bottom: 16px; font-style: italic;">
        💡 Mean = average of all response times (sensitive to outliers) | Median = middle value (typical response time)
      </div>

      <div style="display: flex; gap: 20px; flex-wrap: wrap;">
        <div style="flex: 1; min-width: 200px; padding: 12px; background-color: white; border-radius: 6px; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
          <div style="font-size: 12px; color: #888; margin-bottom: 4px;">📊 Total Messages Analyzed</div>
          <div style="font-size: 24px; font-weight: 600; color: #333;">{activity_info['total_messages']:,}</div>
          {orphaned_warning}
        </div>
        <div style="flex: 1; min-width: 200px; padding: 12px; background-color: white; border-radius: 6px; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
          <div style="font-size: 12px; color: #888; margin-bottom: 4px;">⏱️ Avg Response Time (Mean)</div>
          <div style="font-size: 24px; font-weight: 600; color: #333;">{activity_info['overall_avg_response_minutes']:.1f} min</div>
        </div>
        <div style="flex: 1; min-width: 200px; padding: 12px; background-color: white; border-radius: 6px; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">
          <div style="font-size: 12px; color: #888; margin-bottom: 4px;">⏱️ Median Response Time</div>
          <div style="font-size: 24px; font-weight: 600; color: #333;">{activity_info['overall_median_response_minutes']:.1f} min</div>
        </div>
      </div>
    </div>

    <div class="subtle" style="margin-bottom: 12px; font-style: italic;">
      ℹ️ "Total Messages" includes both sent (by Me) and received (from participant) messages for each conversation{orphaned_explanation}
    </div>

    <table class="meta phone-meta">
      <tr><th>Participant</th><th>Total Messages</th><th>% of Total</th><th>Avg Response (Mean)</th><th>Median Response</th></tr>"""

            # Filter out "Me" from the detailed table since they're already shown in the chart
            for participant, stats in activity_info['participant_stats'].items():
                if participant == "Me":
                    continue
                avg_response = f"{stats['avg_response_time']:.1f} min" if stats['avg_response_time'] > 0 else "N/A"
                median_response = f"{stats['median_response_time']:.1f} min" if stats['median_response_time'] > 0 else "N/A"
                head += f"""
      <tr><td>{html.escape(participant)}</td><td>{stats['message_count']:,}</td><td>{stats['percentage']:.1f}%</td><td>{avg_response}</td><td>{median_response}</td></tr>"""

            head += """
    </table>
  </div>"""

    # Add participant interaction matrix frame if requested
    if show_interaction_matrix:
        interaction_info = get_participant_interaction_matrix(db_path, contact_lookup, tz_from, tz_to)
        if interaction_info:
            head += f"""
  <div class="card">
    <div class="card-header">
      <div class="card-icon info">🔗</div>
      <div>
        <div class="card-title">Participant Interaction Matrix</div>
        <div class="card-subtitle">Who talks to whom most often - {interaction_info['total_interactions']:,} total message interactions</div>
      </div>
    </div>
    <div class="subtle" style="margin-bottom: 16px; font-style: italic;">
      ℹ️ Interaction totals count messages between two people across ALL their shared conversations
    </div>

    <table class="meta phone-meta">
      <tr><th>Interaction Pair</th><th>Total Messages</th><th>%</th><th>Shared Chats</th></tr>"""

            for interaction in interaction_info['interaction_pairs']:
                head += f"""
      <tr><td>{html.escape(interaction['person1'])} ↔ {html.escape(interaction['person2'])}</td><td>{interaction['message_count']:,}</td><td>{interaction['percentage']:.1f}%</td><td>{interaction['chat_count']}</td></tr>"""

            head += """
    </table>
  </div>

  <div class="card">
    <div class="card-header">
      <div class="card-icon warning">💬</div>
      <div>
        <div class="card-title">Most Active Conversations</div>
        <div class="card-subtitle">Top conversations by message count</div>
      </div>
    </div>
    <table class="meta phone-meta">
      <tr><th>Chat</th><th>Participants</th><th>Total Messages</th></tr>"""

            for chat in interaction_info['chat_summary']:
                participants_str = ", ".join(chat['participants'][:4])  # Show first 4 participants
                if len(chat['participants']) > 4:
                    participants_str += f" (+{len(chat['participants']) - 4} more)"

                head += f"""
      <tr><td>{html.escape(chat['display_name'])}</td><td>{html.escape(participants_str)}</td><td>{chat['message_count']:,}</td></tr>"""

            head += """
    </table>
  </div>"""

    # Add communication heatmap frame if requested
    if show_communication_heatmap:
        heatmap_info = get_communication_heatmap_info(db_path, chat_id, tz_from, tz_to, date_from_apple, date_to_apple)
        if heatmap_info:
            head += f"""
  <div class="card">
    <div class="card-header">
      <div class="card-icon warning">🔥</div>
      <div>
        <div class="card-title">Communication Heatmap</div>
        <div class="card-subtitle">Activity patterns by time and day - {heatmap_info['total_messages']:,} messages in {heatmap_info['timezone_used']}</div>
      </div>
    </div>

    <div class="metrics-grid" style="margin-bottom: 20px;">
      <div class="stat-row">
        <span class="stat-label">🔆 Peak Hour</span>
        <span class="stat-value">{heatmap_info['peak_hour']:02d}:00 ({heatmap_info['peak_hour_count']} msgs)</span>
      </div>
      <div class="stat-row">
        <span class="stat-label">🌙 Quietest Hour</span>
        <span class="stat-value">{heatmap_info['quietest_hour']:02d}:00 ({heatmap_info['quietest_hour_count']} msgs)</span>
      </div>
      <div class="stat-row">
        <span class="stat-label">📅 Most Active Day</span>
        <span class="stat-value">{heatmap_info['peak_day']} ({heatmap_info['peak_day_count']} msgs)</span>
      </div>
      <div class="stat-row">
        <span class="stat-label">📅 Quietest Day</span>
        <span class="stat-value">{heatmap_info['quietest_day']} ({heatmap_info['quietest_day_count']} msgs)</span>
      </div>
    </div>

    <h3 style="font-size: 16px; font-weight: 600; margin: 20px 0 10px 0;">Activity by Time Period</h3>
    <table class="meta phone-meta">
      <tr><th>Time Period</th><th>Hours</th><th>Messages</th><th style="width: 30%;">%</th></tr>"""

            for block in heatmap_info['time_blocks']:
                head += f"""
      <tr>
        <td>{block['name']}</td>
        <td>{block['hours']}</td>
        <td>{block['count']:,}</td>
        <td>
          <div style="display: flex; align-items: center; gap: 8px;">
            <div style="flex: 1; background-color: #f0f0f0; border-radius: 3px; height: 18px; overflow: hidden;">
              <div style="width: {block['percentage']:.1f}%; background: linear-gradient(90deg, #007AFF, #5AC8FA); height: 100%;" title="{block['percentage']:.1f}%">
              </div>
            </div>
            <span style="min-width: 45px; text-align: right; font-size: 13px;">{block['percentage']:.1f}%</span>
          </div>
        </td>
      </tr>"""

            head += """
    </table>

    <h3 style="font-size: 16px; font-weight: 600; margin: 20px 0 10px 0;">Hourly Activity (24-hour)</h3>
    <table class="meta phone-meta">
      <tr><th>Hour</th><th>Messages</th><th style="width: 40%;">Activity</th></tr>"""

            max_hourly = max(count for _, count, _ in heatmap_info['hourly_activity'])
            for hour, count, percentage in heatmap_info['hourly_activity']:
                # Calculate percentage width for progress bar
                bar_width = (count / max_hourly * 100) if max_hourly > 0 else 0
                head += f"""
      <tr>
        <td>{hour:02d}:00</td>
        <td>{count:,}</td>
        <td>
          <div class="progress-bar-container">
            <div class="progress-bar-fill" style="width: {bar_width}%;">
              {f'<span class="progress-bar-fill-text">{percentage:.1f}%</span>' if bar_width > 15 else ''}
            </div>
            {f'<span class="progress-bar-text">{percentage:.1f}%</span>' if bar_width <= 15 else ''}
          </div>
        </td>
      </tr>"""

            head += """
    </table>

    <h3 style="font-size: 16px; font-weight: 600; margin: 20px 0 10px 0;">Weekly Activity</h3>
    <table class="meta phone-meta">
      <tr><th>Day</th><th>Messages</th><th style="width: 40%;">Activity</th></tr>"""

            max_daily = max(count for _, count, _ in heatmap_info['daily_activity'])
            for day, count, percentage in heatmap_info['daily_activity']:
                # Calculate percentage width for progress bar
                bar_width = (count / max_daily * 100) if max_daily > 0 else 0
                head += f"""
      <tr>
        <td>{day}</td>
        <td>{count:,}</td>
        <td>
          <div class="progress-bar-container">
            <div class="progress-bar-fill" style="width: {bar_width}%;">
              {f'<span class="progress-bar-fill-text">{percentage:.1f}%</span>' if bar_width > 15 else ''}
            </div>
            {f'<span class="progress-bar-text">{percentage:.1f}%</span>' if bar_width <= 15 else ''}
          </div>
        </td>
      </tr>"""

            head += """
    </table>
  </div>"""

    # Add hash verification frame if requested
    if show_hash_verification and hash_verification_info:
        verified = hash_verification_info['database_hashes']['integrity_verified']
        head += f"""
  <div class="card">
    <div class="card-header">
      <div class="card-icon {'success' if verified else 'warning'}">{'🔒' if verified else '⚠️'}</div>
      <div>
        <div class="card-title">Hash Verification & Data Integrity</div>
        <div class="card-subtitle">Forensic hash verification performed at {hash_verification_info['verification_time']}</div>
      </div>
    </div>

    <div style="background-color: {'#d4edda' if verified else '#f8d7da'}; padding: 15px; border-radius: 8px; margin-bottom: 20px; border-left: 4px solid {'#28a745' if verified else '#dc3545'};">
      <div style="font-weight: 600; font-size: 14px; margin-bottom: 5px;">
        {'✓ INTEGRITY VERIFIED' if verified else '⚠ INTEGRITY CHECK FAILED'}
      </div>
      <div style="font-size: 13px; color: {'#155724' if verified else '#721c24'};">
        {'No changes detected during processing - database integrity confirmed' if verified else 'Database may have been modified during processing'}
      </div>
    </div>

    <table class="meta phone-meta">
      <tr><th>Database File</th><td>{html.escape(hash_verification_info['database_hashes']['pre_processing']['file_path'] if hash_verification_info['database_hashes']['pre_processing'] else str(db_path))}</td></tr>"""

        # Only show file size and modification time if hash info exists
        if hash_verification_info['database_hashes']['pre_processing']:
            head += f"""
      <tr><th>File Size</th><td>{hash_verification_info['database_hashes']['pre_processing']['file_size']:,} bytes</td></tr>
      <tr><th>File Modified</th><td>{html.escape(hash_verification_info['database_hashes']['pre_processing']['modification_time'])}</td></tr>"""

        head += """
    </table>

    <table class="meta phone-meta" style="margin-top: 10px;">
      <tr><th style="background: #ddd; font-weight: bold;">Algorithm</th><th style="background: #ddd; font-weight: bold;">Pre-Processing Hash</th><th style="background: #ddd; font-weight: bold;">Post-Processing Hash</th><th style="background: #ddd; font-weight: bold;">Match</th></tr>"""

        if hash_verification_info['database_hashes']['pre_processing'] and hash_verification_info['database_hashes']['post_processing']:
            pre_hashes = hash_verification_info['database_hashes']['pre_processing']
            post_hashes = hash_verification_info['database_hashes']['post_processing']

            for algorithm in ['md5', 'sha1', 'sha256']:
                pre_hash = pre_hashes[algorithm]
                post_hash = post_hashes[algorithm]
                match = '✓' if pre_hash == post_hash else '✗'
                match_style = 'color: green;' if pre_hash == post_hash else 'color: red; font-weight: bold;'

                head += f"""
      <tr><td style="font-weight: bold;">{algorithm.upper()}</td><td style="font-family: monospace; font-size: 11px;">{pre_hash}</td><td style="font-family: monospace; font-size: 11px;">{post_hash}</td><td style="{match_style}">{match}</td></tr>"""

        head += """
    </table>"""

        # Add AddressBook hash verification if available
        if hash_verification_info['addressbook_hashes']:
            ab_data = hash_verification_info['addressbook_hashes']
            ab_pre = ab_data['pre_processing']
            ab_post = ab_data['post_processing']
            ab_verified = ab_data['integrity_verified']

            head += f"""
    <h3 style="font-size: 16px; font-weight: 600; margin: 30px 0 15px 0;">AddressBook Database Verification</h3>

    <div style="background-color: {'#d4edda' if ab_verified else '#f8d7da'}; padding: 15px; border-radius: 8px; margin-bottom: 20px; border-left: 4px solid {'#28a745' if ab_verified else '#dc3545'};">
      <div style="font-weight: 600; font-size: 14px; margin-bottom: 5px;">
        {'✓ INTEGRITY VERIFIED' if ab_verified else '⚠ INTEGRITY CHECK FAILED'}
      </div>
      <div style="font-size: 13px; color: {'#155724' if ab_verified else '#721c24'};">
        {'No changes detected during processing - AddressBook integrity confirmed' if ab_verified else 'AddressBook may have been modified during processing'}
      </div>
    </div>

    <table class="meta phone-meta">
      <tr><th>AddressBook File</th><td>{html.escape(ab_pre['file_path'] if ab_pre else 'N/A')}</td></tr>
      <tr><th>File Size</th><td>{ab_pre['file_size']:,} bytes</td></tr>
      <tr><th>File Modified</th><td>{html.escape(ab_pre['modification_time']) if ab_pre else 'N/A'}</td></tr>
    </table>

    <table class="meta phone-meta" style="margin-top: 10px;">
      <tr><th style="background: #ddd; font-weight: bold;">Algorithm</th><th style="background: #ddd; font-weight: bold;">Pre-Processing Hash</th><th style="background: #ddd; font-weight: bold;">Post-Processing Hash</th><th style="background: #ddd; font-weight: bold;">Match</th></tr>"""

            if ab_pre and ab_post:
                for algorithm in ['md5', 'sha1', 'sha256']:
                    pre_hash = ab_pre[algorithm]
                    post_hash = ab_post[algorithm]
                    match = '✓' if pre_hash == post_hash else '✗'
                    match_style = 'color: green;' if pre_hash == post_hash else 'color: red; font-weight: bold;'

                    head += f"""
      <tr><td style="font-weight: bold;">{algorithm.upper()}</td><td style="font-family: monospace; font-size: 11px;">{pre_hash}</td><td style="font-family: monospace; font-size: 11px;">{post_hash}</td><td style="{match_style}">{match}</td></tr>"""

            head += """
    </table>"""

        head += """
  </div>"""

    head += """
"""

    report_progress("Generating HTML content...", 45, 100)
    body = ['<div class="messages-section"><h2>💬 Messages</h2>']
    for idx, frame_rows in enumerate(tqdm(frames, desc="Generating HTML frames", unit="frame"), start=1):
        # Report frame progress for GUI (50-70% range)
        frame_progress = 50 + int((idx / len(frames)) * 20) if frames else 50
        report_progress(f"Generating frame {idx}/{len(frames)}", frame_progress, 100)
        frame_html = ["<div id=\"phone-frame\"><div id=\"container\"><div class=\"threaded-chat\">"]
        for row in frame_rows:
            # Clean message text - remove OBJ replacement characters (U+FFFC) used for attachments
            message_text = row["text"] or ""
            # Remove OBJ replacement character (shows as small box with "obj")
            message_text = message_text.replace('\ufffc', '').strip()
            text_html = html.escape(message_text)
            ts_html = html.escape(row["timestamp"] or "")
            cls = "sent" if row["direction"] == "sent" else "received"

            sender_disp = html.escape(row.get('sender_display') or row.get('sender') or '')
            recipient_disp = html.escape(row.get('recipient_display') or row.get('recipient') or '')
            chat_name = html.escape(row.get('chat_display_name') or '')

            # Determine if this is a group chat (more than 2 participants: Me + 1 other)
            is_group_chat = chat_info['participant_count'] > 2

            if row['direction'] == 'sent':
                if is_group_chat and chat_name:
                    who_line = f'Sent: Me → [{chat_name}]'
                else:
                    who_line = f'Sent: Me → {recipient_disp}'
            else:
                if is_group_chat and chat_name:
                    who_line = f'Received: {sender_disp} → [{chat_name}]'
                else:
                    who_line = f'Received: {sender_disp} → Me'
            who_html = f"<div class='meta-line'>{who_line}</div>"

            # Build attachment info if present and enabled
            attachment_html = ""
            if show_attachments and row.get('attachment_filename'):
                attachment_type = row.get('attachment_mime_type') or 'unknown'
                attachment_name = row.get('attachment_transfer_name') or 'Attachment'
                attachment_path_str = row['attachment_filename']
                is_sticker = row.get('attachment_is_sticker', 0)

                # Try to resolve and embed the actual attachment file
                resolved_path = resolve_attachment_path(attachment_path_str, attachments_folder)

                attachment_html = "<div class='attachment-info' style='margin-top: 5px;'>"

                if resolved_path:
                    # File found - try to embed it
                    embedded_html = embed_attachment_in_html(resolved_path, attachment_type)
                    if embedded_html:
                        attachment_html += embedded_html

                    # Add metadata below the embedded content
                    sticker_label = " [Sticker]" if is_sticker else ""
                    attachment_html += f"<div style='font-size: 0.85em; margin-top: 3px; padding: 5px; background-color: rgba(0,0,0,0.05); border-radius: 5px;'>"
                    attachment_html += f"<div style='font-weight: bold;'>📎 {html.escape(attachment_name)}{sticker_label}</div>"
                    attachment_html += f"<div style='font-size: 0.9em; color: #666;'>Type: {html.escape(attachment_type)}</div>"
                    attachment_html += "</div>"
                else:
                    # File not found - show path info only
                    sticker_label = " [Sticker]" if is_sticker else ""
                    attachment_html += f"<div style='font-size: 0.9em; padding: 5px; background-color: rgba(0,0,0,0.05); border-radius: 5px;'>"
                    attachment_html += f"<div style='font-weight: bold;'>📎 {html.escape(attachment_name)}{sticker_label}</div>"
                    attachment_html += f"<div style='font-size: 0.85em; color: #666;'>Type: {html.escape(attachment_type)}</div>"
                    attachment_html += f"<div style='font-size: 0.85em; color: #999; font-style: italic;'>File not found</div>"
                    attachment_html += f"<div style='font-size: 0.75em; color: #999; word-break: break-all;'>Path: {html.escape(attachment_path_str)}</div>"
                    attachment_html += "</div>"

                attachment_html += "</div>"

            frame_html.append(
               f'<div class="bubble {cls}"><b>{who_html}{text_html}</b>{attachment_html}<span class="ts">{ts_html}</span></div>'
            )
        frame_html.append('</div></div></div>')
        frame_html.append(f'<div class="frame-caption">Frame {idx} of {total_frames}</div>')
        frame_html.append('<div class="page-break"></div>')

        body.append("\n".join(frame_html))

    body.append('</div>')  # Close messages-section
    tail = "</div></body>\n</html>\n"  # Close main-container

    # Pagination logic
    if messages_per_page > 0 and len(rows) > messages_per_page:
        # Create a folder from out_path
        if out_path.suffix:
            # If out_path has an extension, create folder with same name minus extension
            output_folder = out_path.parent / out_path.stem
        else:
            # If no extension, use out_path as folder name
            output_folder = out_path

        output_folder.mkdir(parents=True, exist_ok=True)

        # Split rows into chunks
        total_pages = math.ceil(len(rows) / messages_per_page)

        # Generate individual page HTML files
        report_progress("Generating paginated HTML...", 70, 100)
        print(f"Generating {total_pages} paginated HTML files...")
        for page_num in tqdm(range(1, total_pages + 1), desc="Generating pages", unit="page"):
            # Report pagination progress (70-90% range)
            page_progress = 70 + int((page_num / total_pages) * 20)
            report_progress(f"Generating page {page_num}/{total_pages}", page_progress, 100)

            start_idx = (page_num - 1) * messages_per_page
            end_idx = min(start_idx + messages_per_page, len(rows))
            page_rows = rows[start_idx:end_idx]

            # Get date range for this page
            first_msg_date = page_rows[0].get('timestamp', 'N/A') if page_rows else 'N/A'
            last_msg_date = page_rows[-1].get('timestamp', 'N/A') if page_rows else 'N/A'

            # Build page body with pagination nav
            page_body = []

            # Top pagination navigation
            page_body.append(generate_pagination_nav(page_num, total_pages))

            # Build the message frames for this page
            # If per_frame is set, split page_rows into frames; otherwise treat all page_rows as one frame
            if per_frame_val > 0:
                page_frames_count = math.ceil(len(page_rows) / per_frame_val)
                page_frames = [page_rows[i * per_frame_val : (i + 1) * per_frame_val] for i in range(page_frames_count)]
            else:
                page_frames = [page_rows]
                page_frames_count = 1

            # Render each frame on this page
            for frame_idx, frame_rows in enumerate(page_frames, start=1):
                frame_html = ["<div id=\"phone-frame\"><div id=\"container\"><div class=\"threaded-chat\">"]
                for row in frame_rows:
                    # Clean message text - remove OBJ replacement characters (U+FFFC) used for attachments
                    message_text = row["text"] or ""
                    # Remove OBJ replacement character (shows as small box with "obj")
                    message_text = message_text.replace('\ufffc', '').strip()
                    text_html = html.escape(message_text)
                    ts_html = html.escape(row["timestamp"] or "")
                    cls = "sent" if row["direction"] == "sent" else "received"

                    sender_disp = html.escape(row.get('sender_display') or row.get('sender') or '')
                    recipient_disp = html.escape(row.get('recipient_display') or row.get('recipient') or '')
                    chat_name = html.escape(row.get('chat_display_name') or '')

                    # Determine if this is a group chat (more than 2 participants: Me + 1 other)
                    is_group_chat = chat_info['participant_count'] > 2

                    if row['direction'] == 'sent':
                        if is_group_chat and chat_name:
                            who_line = f'Sent: Me → [{chat_name}]'
                        else:
                            who_line = f'Sent: Me → {recipient_disp}'
                    else:
                        if is_group_chat and chat_name:
                            who_line = f'Received: {sender_disp} → [{chat_name}]'
                        else:
                            who_line = f'Received: {sender_disp} → Me'
                    who_html = f"<div class='meta-line'>{who_line}</div>"

                    # Build attachment info if present and enabled
                    attachment_html = ""
                    if show_attachments and row.get('attachment_filename'):
                        attachment_type = row.get('attachment_mime_type') or 'unknown'
                        attachment_name = row.get('attachment_transfer_name') or 'Attachment'
                        attachment_path_str = row['attachment_filename']
                        is_sticker = row.get('attachment_is_sticker', 0)

                        # Try to resolve and embed the actual attachment file
                        resolved_path = resolve_attachment_path(attachment_path_str, attachments_folder)

                        attachment_html = "<div class='attachment-info' style='margin-top: 5px;'>"

                        if resolved_path:
                            # File found - try to embed it
                            embedded_html = embed_attachment_in_html(resolved_path, attachment_type)
                            if embedded_html:
                                attachment_html += embedded_html

                            # Add metadata below the embedded content
                            sticker_label = " [Sticker]" if is_sticker else ""
                            attachment_html += f"<div style='font-size: 0.85em; margin-top: 3px; padding: 5px; background-color: rgba(0,0,0,0.05); border-radius: 5px;'>"
                            attachment_html += f"<div style='font-weight: bold;'>📎 {html.escape(attachment_name)}{sticker_label}</div>"
                            attachment_html += f"<div style='font-size: 0.9em; color: #666;'>Type: {html.escape(attachment_type)}</div>"
                            attachment_html += "</div>"
                        else:
                            # File not found - show path info only
                            sticker_label = " [Sticker]" if is_sticker else ""
                            attachment_html += f"<div style='font-size: 0.9em; padding: 5px; background-color: rgba(0,0,0,0.05); border-radius: 5px;'>"
                            attachment_html += f"<div style='font-weight: bold;'>📎 {html.escape(attachment_name)}{sticker_label}</div>"
                            attachment_html += f"<div style='font-size: 0.85em; color: #666;'>Type: {html.escape(attachment_type)}</div>"
                            attachment_html += f"<div style='font-size: 0.85em; color: #999; font-style: italic;'>File not found</div>"
                            attachment_html += f"<div style='font-size: 0.75em; color: #999; word-break: break-all;'>Path: {html.escape(attachment_path_str)}</div>"
                            attachment_html += "</div>"

                        attachment_html += "</div>"

                    frame_html.append(
                       f'<div class="bubble {cls}"><b>{who_html}{text_html}</b>{attachment_html}<span class="ts">{ts_html}</span></div>'
                    )
                frame_html.append('</div></div></div>')

                # Add frame caption - show frame number if frames are being used
                if page_frames_count > 1:
                    frame_html.append(f'<div class="frame-caption">Frame {frame_idx} of {page_frames_count} (Page {page_num})</div>')
                else:
                    frame_html.append(f'<div class="frame-caption">Messages {start_idx + 1}-{end_idx} of {len(rows)}</div>')

                # Add page break between frames (but not after the last frame)
                if frame_idx < page_frames_count:
                    frame_html.append('<div class="page-break"></div>')

                page_body.append("\n".join(frame_html))

            # Bottom pagination navigation
            page_body.append(generate_pagination_nav(page_num, total_pages))

            # Write page HTML - use CSS-only header (no analysis sections)
            page_html = css_head + "\n".join(page_body) + tail
            page_file = output_folder / f"page_{page_num}.html"
            page_file.write_text(page_html, encoding="utf-8")

        print(f"Generated {total_pages} paginated HTML files in {output_folder}")

        # Create index.html with overview
        index_body = []

        # Overview section
        index_body.append('<div class="header-container">')
        index_body.append(f'<div class="chat-title phone-meta">MESSAGE ARCHIVE OVERVIEW</div>')
        index_body.append(f'<div class="subtle phone-meta">Paginated view - {total_pages} pages, {messages_per_page} messages per page</div>')
        index_body.append('<table class="meta phone-meta">')
        index_body.append(f'<tr><th>Total Messages</th><td>{len(rows)}</td></tr>')
        index_body.append(f'<tr><th>Participants</th><td>{html.escape(participants_display)}</td></tr>')

        # Add date range if available
        if rows:
            first_date = rows[0].get('timestamp', 'N/A')
            last_date = rows[-1].get('timestamp', 'N/A')
            index_body.append(f'<tr><th>Date Range</th><td>{html.escape(first_date)} to {html.escape(last_date)}</td></tr>')

        index_body.append(f'<tr><th>Messages per Page</th><td>{messages_per_page}</td></tr>')
        index_body.append(f'<tr><th>Total Pages</th><td>{total_pages}</td></tr>')
        index_body.append('</table>')
        index_body.append('</div>')

        # Page links section
        index_body.append('<div class="header-container">')
        index_body.append('<div class="chat-title phone-meta">PAGE INDEX</div>')
        index_body.append('<table class="meta phone-meta">')
        index_body.append('<tr><th style="background: #ddd; font-weight: bold;">Page</th><th style="background: #ddd; font-weight: bold;">Messages</th><th style="background: #ddd; font-weight: bold;">Date Range</th><th style="background: #ddd; font-weight: bold;">Link</th></tr>')

        for page_num in range(1, total_pages + 1):
            start_idx = (page_num - 1) * messages_per_page
            end_idx = min(start_idx + messages_per_page, len(rows))
            page_rows = rows[start_idx:end_idx]

            first_msg_date = html.escape(page_rows[0].get('timestamp', 'N/A')) if page_rows else 'N/A'
            last_msg_date = html.escape(page_rows[-1].get('timestamp', 'N/A')) if page_rows else 'N/A'

            index_body.append(f'<tr>')
            index_body.append(f'<td>Page {page_num}</td>')
            index_body.append(f'<td>{start_idx + 1}-{end_idx}</td>')
            index_body.append(f'<td style="font-size: 11px;">{first_msg_date}<br>to<br>{last_msg_date}</td>')
            index_body.append(f'<td><a href="page_{page_num}.html" style="color: #007AFF; text-decoration: none;">View Page {page_num}</a></td>')
            index_body.append(f'</tr>')

        index_body.append('</table>')
        index_body.append('</div>')

        # Create index HTML (reuse head which includes all analysis sections)
        index_html = head + "\n".join(index_body) + tail
        index_file = output_folder / "index.html"
        index_file.write_text(index_html, encoding="utf-8")

        print(f"Index file created: {index_file}")

        # Update out_path to point to index.html for return value
        out_path = index_file
    else:
        # No pagination - write single HTML file (original behavior)
        report_progress("Writing HTML file...", 70, 100)
        html_full = head + "\n".join(body) + tail
        out_path.write_text(html_full, encoding="utf-8")

    report_progress("HTML generation complete", 90, 100)

    # Generate PowerPoint if requested
    if powerpoint_path:
        report_progress("Generating PowerPoint...", 90, 100)
        try:
            # Use powerpoint_frames if available, otherwise fall back to HTML frames
            frames_to_use = powerpoint_frames if powerpoint_frames else frames
            pptx_path = create_powerpoint_from_frames(
                frames_to_use, participants_display, title_display,
                width, font_size, powerpoint_path, chat_info
            )
            print(f"PowerPoint presentation saved to: {pptx_path}")
        except Exception as e:
            print(f"Error creating PowerPoint: {e}")
            print("HTML output was still generated successfully.")

    # No image generation for now - could be added later
    image_paths: List[Path] = []

    # Report fallback contact usage for forensic audit trail
    fallback_usage = {}
    if contact_lookup and contact_lookup.fallback_usage_log:
        print("\n" + "="*60)
        print("FORENSIC NOTICE: Fallback contacts were used")
        print("="*60)
        for phone, count in contact_lookup.fallback_usage_log.items():
            contact_name = contact_lookup.fallback_contacts.get(phone, "Unknown")
            fallback_usage[phone] = {"name": contact_name, "count": count}
            print(f"  {phone} -> '{contact_name}' (used {count} time(s))")
        print("="*60 + "\n")

    # Write contact lookup diagnostic report
    if contact_lookup:
        diagnostic_path = out_path.parent / "contact_lookup_diagnostics.txt"
        contact_lookup.write_diagnostic_report(str(diagnostic_path))
        print(f"\n💡 TIP: Check {diagnostic_path.name} to see which phone numbers couldn't be matched to contacts")

    report_progress("Conversion complete!", 100, 100)
    return out_path, image_paths, fallback_usage


def calculate_optimal_messages_per_frame(rows, phone_w, font_size, max_slide_height=8, chat_info=None):
    """Calculate optimal number of messages per frame based on slide height constraints"""
    if not rows:
        return len(rows), [rows]

    if chat_info is None:
        chat_info = {"participant_count": 2}  # Default to 1-on-1

    # Convert slide height from inches to pixels (approximate)
    max_phone_height_px = int(max_slide_height * 96)  # 96 DPI

    # Calculate message heights
    border_width = 6
    container_padding = 10
    bubble_margin = 6
    bubble_padding_y = 10
    base_phone_height = 100  # Base padding

    # Estimate font heights
    meta_font_size = max(9, int(font_size * 0.79))

    optimal_frames = []
    current_frame = []
    current_height = base_phone_height + border_width * 2

    for row in rows:
        # Estimate message height
        text = row.get('text', '')
        sender_display = row.get('sender_display', '')
        recipient_display = row.get('recipient_display', '')
        chat_name = row.get('chat_display_name', '')

        # Calculate bubble content
        direction = row.get("direction", "received")
        is_group_chat = chat_info['participant_count'] > 2

        if direction == 'sent':
            if is_group_chat and chat_name:
                who_line = f'Sent: Me → [{chat_name}]'
            else:
                who_line = f'Sent: Me → {recipient_display}'
        else:
            if is_group_chat and chat_name:
                who_line = f'Received: {sender_display} → [{chat_name}]'
            else:
                who_line = f'Received: {sender_display} → Me'

        # Estimate wrapped text lines
        message_length = len(text) if text else 0
        if message_length > 100:
            adjusted_max_width = int(phone_w * 0.95)
        elif message_length > 50:
            adjusted_max_width = int(phone_w * 0.90)
        else:
            adjusted_max_width = int(phone_w * 0.85)

        max_chars_per_line = int(adjusted_max_width / (font_size * 0.6))
        estimated_text_lines = max(1, len(text) // max_chars_per_line + (1 if len(text) % max_chars_per_line else 0)) if text else 0

        # Estimate bubble height: who_line + text_lines + timestamp + padding + spacing
        estimated_bubble_height = (
            meta_font_size + 10 +  # who_line + spacing
            (estimated_text_lines * font_size) + (estimated_text_lines * 4) +  # text lines with spacing
            meta_font_size + 12 +  # timestamp + spacing
            bubble_padding_y * 2 +  # bubble padding
            bubble_margin * 3  # bubble margins
        )

        # Check if adding this message would exceed slide height
        projected_height = current_height + estimated_bubble_height

        if projected_height > max_phone_height_px and current_frame:
            # Start new frame
            optimal_frames.append(current_frame)
            current_frame = [row]
            current_height = base_phone_height + border_width * 2 + estimated_bubble_height
        else:
            # Add to current frame
            current_frame.append(row)
            current_height = projected_height

    # Add remaining messages
    if current_frame:
        optimal_frames.append(current_frame)

    # Calculate average messages per frame
    if optimal_frames:
        avg_per_frame = len(rows) // len(optimal_frames)
        return avg_per_frame, optimal_frames
    else:
        return len(rows), [rows]


def create_powerpoint_from_frames(
    frames,
    participants_display,
    title_display,
    phone_w,
    font_size,
    pptx_path,
    chat_info=None,
):
    """Create a PowerPoint presentation with one phone frame per slide"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        raise ImportError(f"PowerPoint generation requires additional packages. Install with: pip install python-pptx pillow\nError: {e}")

    if chat_info is None:
        chat_info = {"participant_count": 2}  # Default to 1-on-1

    # Create presentation
    prs = Presentation()

    # Set slide dimensions to 16:9 (standard)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Colors for phone frame
    sent_color = "#0b93f6"
    recv_color = "#e5e5ea"
    sent_text = "#ffffff"
    recv_text = "#000000"

    # Try to load fonts (bold for message content + emoji support)
    try:
        # Try Windows fonts with bold for content
        title_font = ImageFont.truetype("seguiemj.ttf", 24)  # Segoe UI Emoji for title
        body_font = ImageFont.truetype("seguisb.ttf", font_size)  # Segoe UI Semibold for BOLD content
        meta_font = ImageFont.truetype("segoeui.ttf", max(9, int(font_size * 0.79)))  # Segoe UI for meta
    except:
        try:
            # Fallback to Arial with explicit bold
            title_font = ImageFont.truetype("arial.ttf", 24)
            body_font = ImageFont.truetype("arialbd.ttf", font_size)  # Arial Bold for content
            meta_font = ImageFont.truetype("arial.ttf", max(9, int(font_size * 0.79)))
        except:
            try:
                # Third fallback - regular fonts
                title_font = ImageFont.truetype("arial.ttf", 24)
                body_font = ImageFont.truetype("arial.ttf", font_size)
                meta_font = ImageFont.truetype("arial.ttf", max(9, int(font_size * 0.79)))
            except:
                # Final fallback to default fonts
                title_font = ImageFont.load_default()
                body_font = ImageFont.load_default()
                meta_font = ImageFont.load_default()

    print(f"Generating PowerPoint with {len(frames)} slides...")
    for frame_idx, frame_rows in enumerate(tqdm(frames, desc="Creating PowerPoint slides", unit="slide"), start=1):
        # Add a new slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"{title_display} - Frame {frame_idx}"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True

        # Create phone frame image
        phone_image = create_phone_frame_image(
            frame_rows, phone_w, font_size, sent_color, recv_color,
            sent_text, recv_text, body_font, meta_font, chat_info
        )

        # Save image to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            phone_image.save(tmp_file.name, 'PNG')
            tmp_path = tmp_file.name

        # Add image to slide (centered, larger)
        img_width = Inches(7)  # Larger width
        img_height = Inches(8)  # Taller image
        left = Inches(4.5)  # Center horizontally
        top = Inches(0.8)  # Closer to title

        slide.shapes.add_picture(tmp_path, left, top, img_width, img_height)

        # Participants info removed as requested

        # Clean up temp file
        try:
            Path(tmp_path).unlink()
        except:
            pass

    # Save presentation
    prs.save(str(pptx_path))
    return pptx_path


def create_phone_frame_image(frame_rows, phone_w, font_size, sent_color, recv_color, sent_text, recv_text, body_font, meta_font, chat_info=None):
    """Create an image of a phone frame with messages matching HTML layout"""
    try:
        from PIL import Image, ImageDraw
        import textwrap
    except ImportError:
        raise ImportError("Image generation requires Pillow: pip install pillow")

    if chat_info is None:
        chat_info = {"participant_count": 2}  # Default to 1-on-1

    # Phone dimensions matching HTML design
    border_width = 6
    corner_radius = 36
    phone_inner_width = phone_w

    # Calculate estimated height needed based on message content
    estimated_height = 100  # Base padding
    bubble_padding_y = 10  # Define here for this function
    bubble_margin = 6

    for row in frame_rows:
        text = row.get('text', '')
        text_length = len(text)

        # More accurate height estimation based on wrapping
        if text_length > 100:
            adjusted_max_width = int(phone_w * 0.95)
        elif text_length > 50:
            adjusted_max_width = int(phone_w * 0.90)
        else:
            adjusted_max_width = int(phone_w * 0.85)

        max_chars_per_line = int(adjusted_max_width / (font_size * 0.6))
        estimated_text_lines = max(1, text_length // max_chars_per_line + (1 if text_length % max_chars_per_line else 0)) if text else 0

        # Calculate message bubble height
        meta_font_size = max(9, int(font_size * 0.79))
        estimated_bubble_height = (
            meta_font_size + 10 +  # who_line + spacing
            (estimated_text_lines * font_size) + (estimated_text_lines * 4) +  # text lines
            meta_font_size + 12 +  # timestamp + spacing
            bubble_padding_y * 2 +  # bubble padding
            bubble_margin * 3  # bubble margins
        )
        estimated_height += estimated_bubble_height

    phone_inner_height = max(600, estimated_height)  # Dynamic height with minimum
    total_width = phone_inner_width + border_width * 2
    total_height = phone_inner_height + border_width * 2

    # Padding and spacing matching HTML
    container_padding = 10  # Matches #container padding in HTML
    bubble_margin = 6  # Matches .bubble margin
    bubble_padding_x = 14  # Matches .bubble padding
    bubble_padding_y = 10
    max_bubble_width = int(phone_inner_width * 0.85)  # Wider bubbles (was 75%, now 85%)

    # Create image with background
    img = Image.new('RGB', (total_width, total_height), '#f3f3f3')
    draw = ImageDraw.Draw(img)

    # Draw phone frame with rounded corners
    phone_rect = (0, 0, total_width, total_height)
    try:
        draw.rounded_rectangle(phone_rect, radius=corner_radius, fill='#ffffff', outline='#000000', width=border_width)
    except AttributeError:
        draw.rectangle(phone_rect, fill='#ffffff', outline='#000000', width=border_width)

    # Start drawing messages from top with proper spacing
    y_pos = border_width + container_padding

    for row in frame_rows:
        # Dynamic check - if we're getting close to bottom, expand phone height
        if y_pos > phone_inner_height - 100:
            # Expand phone frame if needed
            phone_inner_height += 200
            total_height = phone_inner_height + border_width * 2
            # Recreate image with new size
            new_img = Image.new('RGB', (total_width, total_height), '#f3f3f3')
            new_draw = ImageDraw.Draw(new_img)
            # Copy existing content
            new_img.paste(img, (0, 0))
            img = new_img
            draw = new_draw
            # Redraw phone frame
            phone_rect = (0, 0, total_width, total_height)
            try:
                draw.rounded_rectangle(phone_rect, radius=corner_radius, fill='#ffffff', outline='#000000', width=border_width)
            except AttributeError:
                draw.rectangle(phone_rect, fill='#ffffff', outline='#000000', width=border_width)

        direction = row.get("direction", "received")
        text = row.get("text", "")
        timestamp = row.get("timestamp", "")
        sender_display = row.get('sender_display', '')
        recipient_display = row.get('recipient_display', '')
        chat_name = row.get('chat_display_name', '')

        # Determine if this is a group chat
        is_group_chat = chat_info['participant_count'] > 2

        # Set bubble properties based on direction
        if direction == 'sent':
            if is_group_chat and chat_name:
                who_line = f'Sent: Me → [{chat_name}]'
            else:
                who_line = f'Sent: Me → {recipient_display}'
            bubble_color = sent_color
            text_color = sent_text
            align_right = True
        else:
            if is_group_chat and chat_name:
                who_line = f'Received: {sender_display} → [{chat_name}]'
            else:
                who_line = f'Received: {sender_display} → Me'
            bubble_color = recv_color
            text_color = recv_text
            align_right = False

        # Calculate initial bubble width and adjust for long messages
        initial_max_width = int(phone_inner_width * 0.85)  # Standard 85% width

        # Check if message is very long and might need wider bubble
        message_length = len(text) if text else 0
        if message_length > 100:  # Long message threshold
            # For long messages, use up to 95% of phone width
            adjusted_max_width = int(phone_inner_width * 0.95)
        elif message_length > 50:  # Medium-long message
            # For medium messages, use 90% of phone width
            adjusted_max_width = int(phone_inner_width * 0.90)
        else:
            # Normal messages use standard width
            adjusted_max_width = initial_max_width

        # Wrap text based on adjusted width
        max_chars_per_line = int(adjusted_max_width / (font_size * 0.6))  # Estimate chars per line
        wrapped_text = textwrap.fill(text, width=max_chars_per_line) if text else ""
        # Don't wrap the who_line to prevent it from moving to content area
        wrapped_who = who_line  # Keep sent/received line as single line

        # Use the adjusted width as our max bubble width for this message
        max_bubble_width = adjusted_max_width

        # Calculate bubble dimensions with proper text wrapping
        lines = [wrapped_who, wrapped_text, timestamp]
        text_lines = []

        for line in lines:
            if line:
                if '\n' in line:  # Handle wrapped text
                    text_lines.extend(line.split('\n'))
                else:
                    text_lines.append(line)

        # Calculate bubble size with proper spacing between sections
        max_line_width = 0
        total_text_height = 0

        for i, line in enumerate(text_lines):
            if line:
                font = meta_font if (i == 0 or i == len(text_lines)-1) else body_font  # First/last lines use meta_font
                try:
                    bbox = draw.textbbox((0, 0), line, font=font)
                    line_width = bbox[2] - bbox[0]
                    line_height = bbox[3] - bbox[1]
                except:
                    line_width = len(line) * (font_size * 0.6)
                    line_height = (font_size if font == body_font else int(font_size * 0.79))

                # Add increased spacing between all elements
                if i == 0:  # After sent/received line
                    line_height += 10  # More space between meta-line and content (was 6)
                elif i == len(text_lines) - 2:  # Before timestamp (last content line)
                    line_height += 12  # More space between content and timestamp (was 8)
                else:
                    line_height += 4  # Increased normal line spacing (was 2)

                max_line_width = max(max_line_width, line_width)
                total_text_height += line_height

        # Set bubble dimensions
        bubble_width = min(max_line_width + bubble_padding_x * 2, max_bubble_width)
        bubble_height = total_text_height + bubble_padding_y * 2

        # Position bubble with margins
        if align_right:
            bubble_x = total_width - border_width - container_padding - bubble_width - bubble_margin
        else:
            bubble_x = border_width + container_padding + bubble_margin

        # Draw bubble with rounded corners
        bubble_rect = (bubble_x, y_pos, bubble_x + bubble_width, y_pos + bubble_height)
        try:
            draw.rounded_rectangle(bubble_rect, radius=18, fill=bubble_color)
        except AttributeError:
            draw.rectangle(bubble_rect, fill=bubble_color)

        # Draw text lines with proper spacing and emoji support
        text_y = y_pos + bubble_padding_y
        for i, line in enumerate(text_lines):
            if line:
                font = meta_font if (i == 0 or i == len(text_lines)-1) else body_font
                try:
                    # Try to draw with emoji support
                    draw.text((bubble_x + bubble_padding_x, text_y), line, font=font, fill=text_color, embedded_color=True)
                except:
                    # Fallback for older Pillow versions
                    draw.text((bubble_x + bubble_padding_x, text_y), line, font=font, fill=text_color)

                try:
                    bbox = draw.textbbox((0, 0), line, font=font)
                    line_height = bbox[3] - bbox[1]
                except:
                    line_height = (font_size if font == body_font else int(font_size * 0.79))

                # Apply increased spacing logic as in calculation
                if i == 0:  # After sent/received line
                    line_height += 10  # More space between meta-line and content (was 6)
                elif i == len(text_lines) - 2:  # Before timestamp (last content line)
                    line_height += 12  # More space between content and timestamp (was 8)
                else:
                    line_height += 4  # Increased normal line spacing (was 2)

                text_y += line_height

        # Move to next bubble position with increased spacing
        y_pos += bubble_height + bubble_margin * 3  # More space between bubbles

    return img

def build_arg_parser():
    ap = argparse.ArgumentParser(description="Convert iMessage SQLite database to phone-style HTML with contact name resolution.")
    ap.add_argument("--db", required=True, help="Path to input SQLite database (chat.db or sms.db)")
    ap.add_argument("--addressbook", help="Path to AddressBook database for contact name resolution")
    ap.add_argument("--out", required=True, help="Path to output HTML")
    ap.add_argument("--width", type=int, default=375, help="Phone frame width in px (default 375)")
    ap.add_argument("--font-size", type=int, default=14, help="Font size for message text in px (default 14)")
    ap.add_argument("--title", default="IMESSAGE MESSAGES", help="Title at top of report")
    ap.add_argument("--chat-id", type=int, help="Filter to a specific chat ID")
    ap.add_argument("--contact-name", help="Friendly name to display alongside the chat")
    ap.add_argument("--tz-from", default="UTC", help="Source timezone for timestamps (default: UTC)")
    ap.add_argument("--tz-to", default=None, help="Destination timezone for timestamps (e.g., America/Chicago)")
    ap.add_argument("--per-frame", type=int, default=0, help="Max messages per phone frame (0 = all messages in one frame)")
    ap.add_argument(
        "--show-tz-in-messages",
        action="store_true",
        help="Append the destination timezone to each timestamp (requires --tz-to)",
    )
    ap.add_argument(
        "--show-phone-numbers",
        action="store_true",
        help="Show both phone number and contact name (e.g., '+16185555555 (John Doe)')",
    )
    ap.add_argument("--powerpoint", "--pptx", help="Generate PowerPoint presentation instead of HTML (specify .pptx file path)")
    ap.add_argument("--date-from", help="Filter messages from this date (YYYY-MM-DD or YYYY-MM-DD HH:MM:SS)")
    ap.add_argument("--date-to", help="Filter messages to this date (YYYY-MM-DD or YYYY-MM-DD HH:MM:SS)")
    ap.add_argument("--show-epoch-info", action="store_true", help="Show epoch timestamp information in the report")
    ap.add_argument("--show-deleted-info", action="store_true", help="Show deleted message analysis in the report")
    ap.add_argument("--show-participant-info", action="store_true", help="Show participant activity and response time analysis in the report")
    ap.add_argument("--show-interaction-matrix", action="store_true", help="Show participant interaction matrix - who talks to whom most often")
    ap.add_argument("--show-communication-heatmap", action="store_true", help="Show communication heatmaps by time of day and day of week in selected timezone")
    ap.add_argument("--show-hash-verification", action="store_true", help="Show hash verification and data integrity information in the report")
    ap.add_argument("--show-attachments", action="store_true", help="Show attachment information (file paths, types) in messages")
    ap.add_argument("--attachments-folder", help="Path to folder containing message attachments (e.g., ~/Library/Messages/Attachments)")
    ap.add_argument("--search", help="Search for messages containing specific text (case-insensitive)")
    ap.add_argument("--messages-per-page", type=int, default=0, help="Split output into multiple pages with this many messages per page (0 = no pagination, default). Creates a folder with index.html and page_N.html files")
    ap.add_argument("--gui", action="store_true", help="Launch a simple GUI for selecting options")
    ap.add_argument("--relations-graph-out", help="Path to write an HTML graph of participant relationships (Cytoscape.js).")
    ap.add_argument("--relations-min-edge", type=int, default=5, help="Only include edges with at least this many messages (default: 5).")
    ap.add_argument("--relations-top-nodes", type=int, default=150, help="Keep at most this many strongest participants (default: 150).")
    ap.add_argument("--relations-exclude-me", action="store_true", help="Exclude 'Me' from the relationship graph.")
    
    return ap

def run_cli(args):
    try:
        out_path, image_paths, fallback_usage = convert_imessage_db(
            args.db,
            args.out,
            addressbook_path=args.addressbook,
            width=args.width,
            font_size=args.font_size,
            title=args.title,
            chat_id=args.chat_id,
            contact_name=args.contact_name,
            tz_from=args.tz_from,
            tz_to=args.tz_to,
            per_frame=args.per_frame,
            show_tz=args.show_tz_in_messages,
            show_phone=args.show_phone_numbers,
            powerpoint_path=args.powerpoint,
            date_from=args.date_from,
            date_to=args.date_to,
            show_epoch_info=args.show_epoch_info,
            show_deleted_info=args.show_deleted_info,
            show_participant_info=args.show_participant_info,
            show_interaction_matrix=args.show_interaction_matrix,
            show_communication_heatmap=args.show_communication_heatmap,
            show_hash_verification=args.show_hash_verification,
            show_attachments=args.show_attachments,
            attachments_folder=args.attachments_folder,
            search_text=args.search,
            messages_per_page=args.messages_per_page,
        )
    except Exception as exc:
        raise SystemExit(str(exc)) from exc
    print(f"Wrote: {out_path}")
    if image_paths:
        print(f"Saved {len(image_paths)} frame image(s) to {image_paths[0].parent}")
    if args.relations_graph_out:
        try:
            output_path = build_interaction_graph_html(
                db_path=args.db,
                out_html=args.relations_graph_out,
                addressbook_path=args.addressbook,
                tz_from=args.tz_from if hasattr(args, "tz_from") else "UTC",
                tz_to=args.tz_to if hasattr(args, "tz_to") else None,
                min_edge_count=args.relations_min_edge,
                top_nodes=args.relations_top_nodes,
                include_me=(not args.relations_exclude_me),
                activity_fn_name="get_participant_activity_pairs",
            )
            print(f"[OK] Wrote relationships graph → {output_path}")
        except Exception as e:
            print(f"[ERROR] Could not build relationships graph: {e}")

class ProgressDialog:
    """Progress dialog for GUI operations"""
    def __init__(self, parent, title="Processing"):
        import tkinter as tk
        from tkinter import ttk

        self.window = tk.Toplevel(parent)
        self.window.title(title)
        self.window.geometry("400x150")
        self.window.resizable(False, False)

        # Make it modal
        self.window.transient(parent)
        self.window.grab_set()

        # Center the window
        self.window.update_idletasks()
        x = (self.window.winfo_screenwidth() // 2) - (400 // 2)
        y = (self.window.winfo_screenheight() // 2) - (150 // 2)
        self.window.geometry(f"400x150+{x}+{y}")

        # Progress label
        self.label = tk.Label(self.window, text="Initializing...", wraplength=380, justify="center")
        self.label.pack(pady=(20, 10))

        # Progress bar
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(self.window, length=350, mode='determinate',
                                           variable=self.progress_var, maximum=100)
        self.progress_bar.pack(pady=10)

        # Percentage label
        self.percent_label = tk.Label(self.window, text="0%")
        self.percent_label.pack(pady=5)

        self.cancelled = False

    def update_progress(self, message, current, total):
        """Update progress bar and message"""
        if total > 0:
            percentage = int((current / total) * 100)
        else:
            percentage = 0

        self.progress_var.set(percentage)
        self.percent_label.config(text=f"{percentage}%")
        self.label.config(text=message)
        self.window.update()

    def close(self):
        """Close the progress dialog"""
        try:
            self.window.grab_release()
            self.window.destroy()
        except:
            pass

def launch_gui():
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk

    root = tk.Tk()
    root.title("iMessage DB to HTML with Contacts")

    db_var = tk.StringVar()
    addressbook_var = tk.StringVar()
    out_var = tk.StringVar()
    width_var = tk.StringVar(value="375")
    font_size_var = tk.StringVar(value="14")
    title_var = tk.StringVar(value="IMESSAGE MESSAGES")
    chat_var = tk.StringVar()
    contact_var = tk.StringVar()
    tz_from_var = tk.StringVar(value="UTC")
    default_dest = DEFAULT_DEST_TZ if DEFAULT_DEST_TZ in TZ_CHOICES else ""
    tz_to_var = tk.StringVar(value=default_dest)
    per_frame_var = tk.StringVar(value="0")
    show_tz_var = tk.BooleanVar(value=False)
    show_phone_var = tk.BooleanVar(value=False)
    show_epoch_var = tk.BooleanVar(value=False)
    show_deleted_var = tk.BooleanVar(value=False)
    show_participant_var = tk.BooleanVar(value=False)
    show_interaction_var = tk.BooleanVar(value=False)
    show_heatmap_var = tk.BooleanVar(value=False)
    show_hash_var = tk.BooleanVar(value=False)
    show_attachments_var = tk.BooleanVar(value=False)
    attachments_folder_var = tk.StringVar()
    powerpoint_var = tk.StringVar()
    date_from_var = tk.StringVar()
    date_to_var = tk.StringVar()
    search_var = tk.StringVar()
    messages_per_page_var = tk.StringVar(value="0")

    frame = tk.Frame(root, padx=10, pady=10)
    frame.pack(fill="both", expand=True)
    frame.columnconfigure(1, weight=1)

    tz_values = [""] + TZ_CHOICES

    chat_combo = ttk.Combobox(frame, textvariable=chat_var, values=["All chats"], state="readonly")
    tz_from_combo = ttk.Combobox(frame, textvariable=tz_from_var, values=tz_values, state="readonly")
    tz_to_combo = ttk.Combobox(frame, textvariable=tz_to_var, values=tz_values, state="readonly")


    if "UTC" in tz_values:
        tz_from_combo.set("UTC")
    else:
        tz_from_combo.current(0)
    if default_dest:
        tz_to_combo.set(default_dest)
    else:
        tz_to_combo.current(0)

    last_loaded_db = {"path": None, "addressbook_path": None}

    build_relations_var = tk.BooleanVar(value=False)
    min_edge_var = tk.IntVar(value=5)

    def on_toggle_relations():
    # enable/disable the spinbox based on checkbox
        state = "normal" if build_relations_var.get() else "disabled"
        try:
            min_edge_spin.config(state=state)
        except Exception:
            pass

    def populate_date_range(force=False):
        """Populate the date_from and date_to fields with message timeframe using destination timezone"""
        path = db_var.get().strip()

        if not path:
            return

        try:
            # Get current chat selection
            chat_value = chat_var.get().strip()
            chat_id_value = None

            if chat_value and chat_value != "All chats":
                try:
                    chat_id_value = int(chat_value.split(":")[0])
                except:
                    pass

            # Get timezone settings from GUI
            tz_from_value = tz_from_var.get().strip() or "UTC"
            tz_to_value = tz_to_var.get().strip() or None

            # Get date range with timezone conversion
            date_from, date_to = get_message_date_range(path, chat_id_value, tz_from_value, tz_to_value)

            # Only update if we got valid dates and fields are empty or forced
            if date_from and date_to:
                if force or not date_from_var.get().strip():
                    date_from_var.set(date_from)
                if force or not date_to_var.get().strip():
                    date_to_var.set(date_to)

        except Exception as e:
            print(f"Warning: Could not populate date range: {e}")

    def refresh_chats(force=False, show_error=False):
        path = db_var.get().strip()
        addressbook_path = addressbook_var.get().strip()

        if not path:
            chat_combo["values"] = ["All chats"]
            chat_var.set("All chats")
            last_loaded_db["path"] = None
            last_loaded_db["addressbook_path"] = None
            return

        # Check if both DB and AddressBook are unchanged
        if not force and path == last_loaded_db["path"] and addressbook_path == last_loaded_db["addressbook_path"]:
            return

        # Initialize contact lookup for chat name enhancement
        contact_lookup = ContactLookup(addressbook_path if addressbook_path else None)

        chat_ids, error = collect_chat_ids(path, contact_lookup)
        values = ["All chats"] + chat_ids if chat_ids else ["All chats"]
        chat_combo.configure(values=values)

        current_value = chat_var.get()
        # Always reset to "All chats" when loading a new database to avoid chat ID mismatches
        if force or current_value not in values:
            chat_var.set(values[0])
            chat_combo.set(values[0])

        if error:
            if show_error:
                messagebox.showerror("Chat IDs", f"Could not read chat IDs:\n{error}")
            last_loaded_db["path"] = None
            last_loaded_db["addressbook_path"] = None
        else:
            last_loaded_db["path"] = path
            last_loaded_db["addressbook_path"] = addressbook_path
            # Auto-populate date range when database is successfully loaded
            populate_date_range()


    def browse_db():
        path = filedialog.askopenfilename(
            title="Select SQLite database file",
            filetypes=[("SQLite Database", "*.db"), ("All Files", "*.*")],
        )
        if path:
            db_var.set(path)
            # Clear the cache to force reload
            last_loaded_db["path"] = None
            last_loaded_db["addressbook_path"] = None
            # Clear date filters to avoid invalid ranges for new database
            date_from_var.set("")
            date_to_var.set("")
            refresh_chats(force=True, show_error=True)

    def browse_addressbook():
        path = filedialog.askopenfilename(
            title="Select AddressBook database file",
            filetypes=[("IOS Address Book Database", "*.sqlitedb"),("Itunes AddressBook Database", "*.abcddb"), ("SQLite Database", "*.db"), ("All Files", "*.*")],
        )
        if path:
            addressbook_var.set(path)
            # Clear the cache to force reload
            last_loaded_db["path"] = None
            last_loaded_db["addressbook_path"] = None
            # Force refresh of chats to reload contact lookup with new AddressBook
            if db_var.get().strip():  # Only refresh if a database is selected
                refresh_chats(force=True, show_error=False)

    def browse_attachments_folder():
        path = filedialog.askdirectory(
            title="Select Attachments Folder (e.g., ~/Library/Messages/Attachments)",
        )
        if path:
            attachments_folder_var.set(path)

    def show_help():
        """Display README.md in a new window"""
        help_window = tk.Toplevel(root)
        help_window.title("iMessage Database Converter - Help")
        help_window.geometry("800x600")

        # Create frame with scrollbar
        help_frame = tk.Frame(help_window)
        help_frame.pack(fill="both", expand=True, padx=10, pady=10)

        # Add scrollbar
        scrollbar = tk.Scrollbar(help_frame)
        scrollbar.pack(side="right", fill="y")

        # Add text widget
        help_text = tk.Text(help_frame, wrap="word", yscrollcommand=scrollbar.set, font=("Courier", 10))
        help_text.pack(side="left", fill="both", expand=True)
        scrollbar.config(command=help_text.yview)

        # Read and display README.md
        readme_path = Path(__file__).parent / "README.md"
        try:
            if readme_path.exists():
                readme_content = readme_path.read_text(encoding='utf-8')
                help_text.insert("1.0", readme_content)
            else:
                help_text.insert("1.0", "README.md file not found.\n\nFor help, visit: https://github.com/anthropics/claude-code")
        except Exception as e:
            help_text.insert("1.0", f"Error loading README.md: {e}")

        help_text.config(state="disabled")  # Make read-only

        # Add close button
        close_btn = tk.Button(help_window, text="Close", command=help_window.destroy)
        close_btn.pack(pady=5)

    def browse_out():
        path = filedialog.asksaveasfilename(
            title="Save HTML as",
            defaultextension=".html",
            filetypes=[("HTML Files", "*.html"), ("All Files", "*.*")],
        )
        if path:
            out_var.set(path)


   

    def run_conversion():
        # Always refresh chats before conversion to ensure contact lookup is current
        refresh_chats(force=True)

        db_path = db_var.get().strip()
        addressbook_path = addressbook_var.get().strip()
        out_path = out_var.get().strip()
        if not db_path:
            messagebox.showerror("Missing Database", "Please choose an input database file.")
            return
        if not out_path:
            messagebox.showerror("Missing Output", "Please choose where to save the HTML file.")
            return
        try:
            width_value = int(width_var.get().strip() or 375)
        except ValueError:
            messagebox.showerror("Invalid Width", "Phone width must be a whole number.")
            return
        try:
            font_size_value = int(font_size_var.get().strip() or 14)
            if font_size_value < 8 or font_size_value > 32:
                raise ValueError("Font size must be between 8 and 32")
        except ValueError:
            messagebox.showerror("Invalid Font Size", "Font size must be a whole number between 8 and 32.")
            return
        try:
            per_frame_value = int(per_frame_var.get().strip() or 0)
        except ValueError:
            messagebox.showerror("Invalid Messages Per Frame", "Messages per frame must be a whole number.")
            return
        try:
            messages_per_page_value = int(messages_per_page_var.get().strip() or 0)
        except ValueError:
            messagebox.showerror("Invalid Messages Per Page", "Messages per page must be a whole number.")
            return

        title_value = title_var.get().strip() or "IMESSAGE MESSAGES"
        chat_value = chat_var.get().strip()
        if chat_value == "All chats" or not chat_value:
            chat_id_value = None
        else:
            # Extract chat ID from the display string
            try:
                chat_id_value = int(chat_value.split(":")[0])
            except:
                chat_id_value = None

        contact_value = contact_var.get().strip() or None
        tz_from_value = tz_from_var.get().strip() or None
        tz_to_value = tz_to_var.get().strip() or None
        show_tz_value = bool(show_tz_var.get())
        show_phone_value = bool(show_phone_var.get())
        show_epoch_value = bool(show_epoch_var.get())
        show_deleted_value = bool(show_deleted_var.get())
        show_participant_value = bool(show_participant_var.get())
        show_interaction_value = bool(show_interaction_var.get())
        show_heatmap_value = bool(show_heatmap_var.get())
        show_hash_value = bool(show_hash_var.get())
        show_attachments_value = bool(show_attachments_var.get())
        search_value = search_var.get().strip() or None
        build_relations_value = bool(build_relations_var.get())
        min_edge_value = max(1, int(min_edge_var.get() or 1))

        if show_tz_value and not tz_to_value:
            messagebox.showerror(
                "Destination Timezone",
                "Select a destination timezone before including it in message times.",
            )
            return

        # Create progress dialog
        progress_dialog = ProgressDialog(root, "Converting Database")

        # Storage for results and exceptions from thread
        result_container = {"html_path": None, "image_paths": None, "fallback_usage": None, "relations_path": None,"error": None}

        if build_relations_value:
            rel_out = str(Path(out_path).with_name(Path(out_path).stem + "_relations.html"))
            try:
                build_interaction_graph_html(
                    db_path=db_path,
                    out_html=rel_out,
                    addressbook_path=addressbook_path or None,
                    tz_from=tz_from_value or "UTC",
                    tz_to=tz_to_value or None,
                    min_edge_count=min_edge_value,
                    top_nodes=150,         # tweak later if you want
                    include_me=True,
                    activity_fn_name="get_participant_activity_pairs",
                )
                result_container["relations_path"] = rel_out
            except Exception as e:
                # Don't fail overall conversion if graph generation hiccups
                print(f"Relationship graph generation failed: {e}")

        
        
        def progress_callback(message, current, total):
            """Thread-safe progress callback"""
            try:
                progress_dialog.update_progress(message, current, total)
            except:
                pass  # Dialog may have been closed

        def run_conversion_thread():
            """Run conversion in background thread"""
            try:
                html_path, image_paths, fallback_usage = convert_imessage_db(
                    db_path,
                    out_path,
                    addressbook_path=addressbook_path or None,
                    width=width_value,
                    font_size=font_size_value,
                    title=title_value,
                    chat_id=chat_id_value,
                    contact_name=contact_value,
                    tz_from=tz_from_value,
                    tz_to=tz_to_value,
                    per_frame=per_frame_value,
                    show_tz=show_tz_value,
                    show_phone=show_phone_value,
                    powerpoint_path=powerpoint_var.get() or None,
                    date_from=date_from_var.get().strip() or None,
                    date_to=date_to_var.get().strip() or None,
                    show_epoch_info=show_epoch_value,
                    show_deleted_info=show_deleted_value,
                    show_participant_info=show_participant_value,
                    show_interaction_matrix=show_interaction_value,
                    show_communication_heatmap=show_heatmap_value,
                    show_hash_verification=show_hash_value,
                    show_attachments=show_attachments_value,
                    attachments_folder=attachments_folder_var.get().strip() or None,
                    search_text=search_value,
                    messages_per_page=messages_per_page_value,
                    progress_callback=progress_callback,
                )
                result_container["html_path"] = html_path
                result_container["image_paths"] = image_paths
                result_container["fallback_usage"] = fallback_usage
                
            except Exception as exc:
                import traceback
                error_details = f"{str(exc)}\n\nTraceback:\n{traceback.format_exc()}"
                result_container["error"] = error_details

        # Start conversion thread
        conversion_thread = threading.Thread(target=run_conversion_thread, daemon=True)
        conversion_thread.start()

        # Wait for thread to complete (check every 100ms)
        def check_thread():
            if conversion_thread.is_alive():
                root.after(100, check_thread)
            else:
                # Thread finished - close progress dialog and show result
                progress_dialog.close()

                if result_container["error"]:
                    messagebox.showerror("Conversion Failed", result_container["error"])
                    return

                html_path = result_container["html_path"]
                image_paths = result_container["image_paths"]
                fallback_usage = result_container["fallback_usage"]
                
                relations_path = result_container.get("relations_path")
                msg = f"Saved HTML to:\n{html_path}"
                if relations_path:
                    msg += f"\n\nSaved relationships graph to:\n{relations_path}"
                messagebox.showinfo("Conversion Complete", msg)


                # Show success message
                message = f"HTML written to:\n{html_path}"
                if image_paths:
                    message += f"\n\nSaved {len(image_paths)} frame image(s) in:\n{image_paths[0].parent}"

                # Show forensic warning if fallback contacts were used
                if fallback_usage:
                    fallback_details = "\n".join([
                        f"  {phone} -> '{info['name']}' (used {info['count']} time(s))"
                        for phone, info in fallback_usage.items()
                    ])
                    messagebox.showwarning(
                        "FORENSIC NOTICE",
                        f"Fallback contacts were used:\n\n{fallback_details}\n\n" +
                        "These contact names came from contacts_fallback.txt, " +
                        "not from the AddressBook database."
                    )

                messagebox.showinfo("Done", message)

        # Start checking for thread completion
        check_thread()

    row = 0
    tk.Label(frame, text="Database file:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    db_entry = tk.Entry(frame, textvariable=db_var)
    db_entry.grid(row=row, column=1, sticky="we", padx=5, pady=4)
    tk.Button(frame, text="Browse", command=browse_db).grid(row=row, column=2, padx=5, pady=4)

    row += 1
    tk.Label(frame, text="AddressBook file:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    addressbook_entry = tk.Entry(frame, textvariable=addressbook_var)
    addressbook_entry.grid(row=row, column=1, sticky="we", padx=5, pady=4)
    tk.Button(frame, text="Browse", command=browse_addressbook).grid(row=row, column=2, padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Attachments Folder (optional):").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    attachments_folder_entry = tk.Entry(frame, textvariable=attachments_folder_var)
    attachments_folder_entry.grid(row=row, column=1, sticky="we", padx=5, pady=4)
    tk.Button(frame, text="Browse", command=browse_attachments_folder).grid(row=row, column=2, padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Output HTML:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    out_entry = tk.Entry(frame, textvariable=out_var)
    out_entry.grid(row=row, column=1, sticky="we", padx=5, pady=4)
    tk.Button(frame, text="Browse", command=browse_out).grid(row=row, column=2, padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Phone width:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=width_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Font size:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=font_size_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Report title:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=title_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Chat:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    chat_combo.grid(row=row, column=1, sticky="we", padx=5, pady=4)
    tk.Button(frame, text="Refresh", command=lambda: refresh_chats(force=True, show_error=True)).grid(row=row, column=2, padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Contact name:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=contact_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Search text:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=search_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Source timezone:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tz_from_combo.grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Destination timezone:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tz_to_combo.grid(row=row, column=1, sticky="we", padx=5, pady=4)

    # Add Check All / Uncheck All button
    row += 1

    def toggle_all_checkboxes():
        """Toggle all checkboxes on/off"""
        # Check if any checkbox is currently checked
        any_checked = (show_tz_var.get() or show_phone_var.get() or show_epoch_var.get() or
                      show_deleted_var.get() or show_participant_var.get() or show_interaction_var.get() or
                      show_heatmap_var.get() or show_hash_var.get() or show_attachments_var.get() or
                      build_relations_var.get())

        # If any are checked, uncheck all. Otherwise, check all.
        new_value = not any_checked

        show_tz_var.set(new_value)
        show_phone_var.set(new_value)
        show_epoch_var.set(new_value)
        show_deleted_var.set(new_value)
        show_participant_var.set(new_value)
        show_interaction_var.set(new_value)
        show_heatmap_var.set(new_value)
        show_hash_var.set(new_value)
        show_attachments_var.set(new_value)
        build_relations_var.set(new_value)

        # Update the relationships spinbox state
        on_toggle_relations()

    tk.Button(frame, text="Check All / Uncheck All", command=toggle_all_checkboxes).grid(
        row=row, column=1, sticky="w", padx=5, pady=(10, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show destination timezone after message times", variable=show_tz_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show phone numbers with contact names", variable=show_phone_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show epoch timestamp information in report", variable=show_epoch_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show deleted message analysis in report", variable=show_deleted_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show participant activity analysis in report", variable=show_participant_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show participant interaction matrix (who talks to whom)", variable=show_interaction_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )
    
    row += 1
   # Wrap in a frame so the checkbox and spinbox sit on one line
    rel_frame = tk.Frame(frame)
    rel_frame.grid(row=row, column=1, sticky="w", padx=5, pady=(0, 4))

    tk.Checkbutton(
        rel_frame,
        text="Create relationships graph (HTML)",
        variable=build_relations_var,
        command=on_toggle_relations
    ).pack(side="left")

    # Spacer
    tk.Label(rel_frame, text="   ").pack(side="left")

    # Label + Spinbox for min edge count
    tk.Label(rel_frame, text="Min edge count:").pack(side="left")

    min_edge_spin = tk.Spinbox(
        rel_frame,
        from_=1, to=9999,
        width=5,
        textvariable=min_edge_var,
        state="disabled"  # becomes enabled when the checkbox is checked
    )
    min_edge_spin.pack(side="left")


    row += 1
    tk.Checkbutton(frame, text="Show communication heatmap (activity by time/day)", variable=show_heatmap_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show hash verification for data integrity", variable=show_hash_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Checkbutton(frame, text="Show attachments (file paths, types, stickers)", variable=show_attachments_var).grid(
        row=row, column=1, sticky="w", padx=5, pady=(0, 4)
    )

    row += 1
    tk.Label(frame, text="PowerPoint Output (optional):").grid(row=row, column=0, sticky="w", padx=5, pady=4)
    pptx_frame = tk.Frame(frame)
    pptx_frame.grid(row=row, column=1, sticky="we", padx=5, pady=4)
    pptx_frame.columnconfigure(0, weight=1)
    tk.Entry(pptx_frame, textvariable=powerpoint_var).grid(row=0, column=0, sticky="we")
    tk.Button(pptx_frame, text="Browse", command=lambda: powerpoint_var.set(
        filedialog.asksaveasfilename(title="Save PowerPoint as", defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
    )).grid(row=0, column=1, padx=(5, 0))

    row += 1
    tk.Label(frame, text="Messages per frame:").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=per_frame_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Messages per page (0=no pagination):").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    messages_per_page_entry = tk.Entry(frame, textvariable=messages_per_page_var)
    messages_per_page_entry.grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="ℹ Pagination creates a folder with index.html", font=("TkDefaultFont", 8), fg="gray").grid(row=row, column=0, columnspan=3, sticky="w", padx=5, pady=(0, 4))

    # Filtering options separator
    row += 1
    tk.Label(frame, text="--- Filtering Options ---", font=("TkDefaultFont", 9, "italic")).grid(row=row, column=0, columnspan=3, pady=(10, 5))

    row += 1
    tk.Label(frame, text="Date from (YYYY-MM-DD):").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=date_from_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)

    row += 1
    tk.Label(frame, text="Date to (YYYY-MM-DD):").grid(row=row, column=0, sticky="e", padx=5, pady=4)
    tk.Entry(frame, textvariable=date_to_var).grid(row=row, column=1, sticky="we", padx=5, pady=4)
    tk.Button(frame, text="Auto-fill", command=lambda: populate_date_range(force=True)).grid(row=row, column=2, padx=5, pady=4)

    row += 1
    tk.Label(frame, text="ℹ Dates shown in destination timezone", font=("TkDefaultFont", 8), fg="gray").grid(row=row, column=0, columnspan=3, sticky="w", padx=5, pady=(0, 4))


    row += 1
    button_frame = tk.Frame(frame)
    button_frame.grid(row=row, column=0, columnspan=3, pady=(10, 0))
    tk.Button(button_frame, text="Help", command=show_help, width=10).pack(side="left", padx=5)
    tk.Button(button_frame, text="Convert", command=run_conversion, width=15).pack(side="left", padx=5)

    # Set up event bindings
    db_entry.bind("<FocusOut>", lambda _event: refresh_chats(force=True))
    chat_combo.bind("<<ComboboxSelected>>", lambda _event: populate_date_range())
    tz_from_combo.bind("<<ComboboxSelected>>", lambda _event: populate_date_range())
    tz_to_combo.bind("<<ComboboxSelected>>", lambda _event: populate_date_range())

    root.mainloop()
    
def build_interaction_graph_html(
    db_path: str,
    out_html: str,
    *,
    addressbook_path: str | None = None,
    tz_from: str = "UTC",
    tz_to: str | None = None,
    min_edge_count: int = 5,
    top_nodes: int = 9999,
    include_me: bool = True,
    activity_fn_name: "get_participant_activity_pairs" | None = None,   # <- set this to YOUR function name
    debug: bool = False,                    # <- writes a sidecar CSV + pills
) -> str:
    """
    Create a self-contained HTML graph of participant relationships by merging:
      • Participant Interaction Matrix (pairwise totals)
      • Participant Activity (you provide the function name)

    Nodes = participants. Edge thickness = total messages. Edge label = count.
    """

    from pathlib import Path
    import json, csv

    # ---------- helpers ----------
    def norm(name: str | None) -> str:
        n = (name or "").strip()
        if not n:
            return "Unknown"
        nl = n.lower()
        if nl in {"me", "myself"}:
            return "Me"
        return n

    def add_edge(agg: dict, a: str, b: str, count: int):
        if not a or not b or a == b:
            return
        key = tuple(sorted((a, b), key=lambda s: s.lower()))  # undirected
        agg[key] = agg.get(key, 0) + int(count or 0)

    def call_optional(fn_name: str, *args, **kwargs):
        fn = globals().get(fn_name)
        if not callable(fn):
            return None, f"Function '{fn_name}' not found."
        try:
            return fn(*args, **kwargs), None
        except TypeError:
            # try positional fallback
            try:
                return fn(*args), None
            except Exception as e:
                return None, f"{fn_name} call failed (positional): {e}"
        except Exception as e:
            return None, f"{fn_name} call failed: {e}"

    def parse_activity_rows(rows):
        """
        Accepts flexible shapes:
          • {'sender','recipient','message_count'}
          • {'person1','person2','message_count'}
          • {'pair':('A','B'), 'count':...}
          • Per-message rows: {'sender','recipient'} (no count) -> aggregate as 1 each row
        Returns dict{(a,b):count}
        """
        out = {}
        for r in rows:
            # try structured first
            a = r.get("sender") or r.get("person1")
            b = r.get("recipient") or r.get("person2")
            if (a is None or b is None) and isinstance(r.get("pair"), (list, tuple)) and len(r["pair"]) >= 2:
                a, b = r["pair"][0], r["pair"][1]
            cnt = r.get("message_count")
            if cnt is None:
                cnt = r.get("count")
            # if still None, treat as 1 (per-message row)
            if cnt is None:
                cnt = 1
            a = norm(a); b = norm(b)
            add_edge(out, a, b, cnt)
        return out

    lookup = ContactLookup(addressbook_path) if addressbook_path else None

    # ---------- 0) Get true total message count from database ----------
    import sqlite3
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM message")
    db_total_messages = cursor.fetchone()[0]
    conn.close()

    # ---------- 1) Interaction Matrix ----------
    inter_edges = {}
    inter = get_participant_interaction_matrix(
        db_path=db_path, contact_lookup=lookup, tz_from=tz_from, tz_to=tz_to
    )
    if inter and inter.get("interaction_pairs"):
        for pair in inter["interaction_pairs"]:
            a = norm(pair.get("person1"))
            b = norm(pair.get("person2"))
            cnt = int(pair.get("message_count", 0))
            add_edge(inter_edges, a, b, cnt)
    inter_edge_count = sum(inter_edges.values())
    inter_pairs = len(inter_edges)

    # ---------- 2) Participant Activity (your function) ----------
    activity_edges = {}
    activity_pairs = 0
    activity_edge_count = 0
    act_err = None

    if activity_fn_name:
        rows, act_err = call_optional(
            activity_fn_name,
            db_path=db_path, contact_lookup=lookup, tz_from=tz_from, tz_to=tz_to
        )
        if rows:
            activity_edges = parse_activity_rows(rows)
            activity_pairs = len(activity_edges)
            activity_edge_count = sum(activity_edges.values())

    # ---------- 3) Merge ----------
    # Use MAX instead of SUM to avoid double-counting the same messages
    # inter_edges = all messages, activity_edges = back-and-forth only
    # We want the total message count, not the sum of two counting methods
    merged = {}
    for d in (inter_edges, activity_edges):
        for k, v in d.items():
            merged[k] = max(merged.get(k, 0), int(v or 0))

    # ---------- 4) Filter (include_me + min_edge) ----------
    # Apply include_me filter but keep min_edge low for initial pass
    filtered = {}
    min_edge = max(1, int(min_edge_count or 1))
    for (a, b), cnt in merged.items():
        if not include_me and (a.lower() == "me" or b.lower() == "me"):
            continue
        # Lower the initial threshold to 1 so we have all edges available in JS
        if cnt >= 1:
            filtered[(a, b)] = cnt

    if not filtered:
        raise RuntimeError(
            "No edges after merging/filtering. Tips: set Min Edge = 1, include 'Me', enlarge Top Nodes."
        )

    # ---------- 5) Node weights + trim to Top N ----------
    node_weight = {}
    for (a, b), cnt in filtered.items():
        node_weight[a] = node_weight.get(a, 0) + cnt
        node_weight[b] = node_weight.get(b, 0) + cnt

    keep_nodes = set(
        n for n, _w in sorted(node_weight.items(), key=lambda x: x[1], reverse=True)[:max(1, int(top_nodes or 1))]
    )
    trimmed = { (a,b):cnt for (a,b),cnt in filtered.items() if a in keep_nodes and b in keep_nodes }

    # recompute weights after trim (for displayed relationships)
    node_weight = {}
    for (a, b), cnt in trimmed.items():
        node_weight[a] = node_weight.get(a, 0) + cnt
        node_weight[b] = node_weight.get(b, 0) + cnt

    # Calculate true total weight for each node from merged (before filtering)
    # This represents ALL messages each person participated in
    node_weight_all = {}
    for (a, b), cnt in merged.items():
        node_weight_all[a] = node_weight_all.get(a, 0) + cnt
        node_weight_all[b] = node_weight_all.get(b, 0) + cnt

    # Special case: "Me" should always show the true database total
    # because "Me" is involved in every single message
    for node in node_weight_all.keys():
        if node.lower() == "me":
            node_weight_all[node] = db_total_messages

    # ---------- 6) Scaling ----------
    if node_weight:
        min_w = min(node_weight.values()); max_w = max(node_weight.values())
    else:
        min_w = max_w = 1

    def scale_node(w):
        if max_w == min_w: return 36
        return int(24 + (w - min_w) * (80 - 24) / (max_w - min_w))  # 24..80

    edge_vals = list(trimmed.values()) or [1]
    emin, emax = min(edge_vals), max(edge_vals)
    slider_max = max(100, emax)  # Set slider max to at least 100 or the max edge count

    def scale_edge(w):
        if emax == emin: return 3
        return 1 + (w - emin) * (10 - 1) / (emax - emin)           # 1..10

    # Build nodes with color coding
    nodes = []
    for n in sorted(keep_nodes):
        weight = node_weight.get(n, 0)
        is_me = n.lower() == "me"
        # Color gradient: more active = darker blue, "Me" = green
        if is_me:
            color = "#10b981"  # green for "Me"
        else:
            # Gradient from light blue to dark blue based on activity
            if max_w > min_w:
                intensity = (weight - min_w) / (max_w - min_w)
                # Map from #93c5fd (light blue) to #1e40af (dark blue)
                r = int(147 - intensity * (147 - 30))
                g = int(197 - intensity * (197 - 64))
                b = int(253 - intensity * (253 - 175))
                color = f"#{r:02x}{g:02x}{b:02x}"
            else:
                color = "#0ea5e9"  # default sky blue

        nodes.append({
            "data": {
                "id": n,
                "label": n,
                "weight": weight,  # Weight based on visible relationships
                "weight_all": node_weight_all.get(n, weight),  # True total including filtered relationships
                "size": scale_node(weight),
                "color": color,
                "is_me": is_me
            }
        })

    # Build edges with separate counts for matrix, activity, and merged
    cy_edges = []
    for (a, b), cnt in sorted(trimmed.items(), key=lambda kv: (-kv[1], kv[0][0].lower(), kv[0][1].lower())):
        count_matrix = inter_edges.get((a, b), 0)
        count_activity = activity_edges.get((a, b), 0)
        count_merged = cnt
        cy_edges.append({
            "data": {
                "id": f"{a}__{b}",
                "source": a,
                "target": b,
                "count_matrix": count_matrix,
                "count_activity": count_activity,
                "count_merged": count_merged,
                "count": count_merged,  # default display
                "width": scale_edge(cnt),
                "label": str(count_merged)  # default label
            }
        })


    total_edges = len(cy_edges); total_nodes = len(nodes)

    # ---------- 6.5) Statistics for panel ----------
    # Top 10 relationships by merged count
    top_relationships = sorted(trimmed.items(), key=lambda x: x[1], reverse=True)[:20]
    top_relationships_data = [
        {"person1": a, "person2": b, "count": cnt}
        for (a, b), cnt in top_relationships
    ]

    # Most active participant (excluding "Me")
    most_active = None
    most_active_count = 0
    for node_id, weight in node_weight.items():
        if node_id.lower() != "me" and weight > most_active_count:
            most_active = node_id
            most_active_count = weight

    # Calculate statistics
    # Use db_total_messages for the forensically accurate total (includes ALL messages, even group chats)
    # Use trimmed (after filtering) for displayed relationships
    total_messages_displayed = sum(trimmed.values())
    avg_messages_per_relationship = total_messages_displayed / len(trimmed) if trimmed else 0

    stats_data = {
        "top_relationships": top_relationships_data,
        "most_active": most_active or "N/A",
        "most_active_count": most_active_count,
        "total_messages": db_total_messages,  # Forensically accurate total from database
        "total_messages_displayed": total_messages_displayed,
        "avg_per_relationship": round(avg_messages_per_relationship, 1),
        "total_pairs": len(trimmed)
    }

    # ---------- 7) Debug sidecar ----------
    if debug:
        dbg_path = str(Path(out_html).with_suffix("")) + "_relations_debug.csv"
        with open(dbg_path, "w", newline="", encoding="utf-8") as f:
            w = csv.writer(f)
            w.writerow(["source", "target", "count_interaction", "count_activity", "count_merged", "kept_after_filters"])
            all_pairs = set(inter_edges.keys()) | set(activity_edges.keys()) | set(merged.keys())
            for a, b in sorted(all_pairs, key=lambda p: (p[0].lower(), p[1].lower())):
                ci = inter_edges.get((a, b), 0)
                ca = activity_edges.get((a, b), 0)
                cm = merged.get((a, b), 0)
                kept = trimmed.get((a, b), 0)
                w.writerow([a, b, ci, ca, cm, kept])

    # ---------- 8) HTML ----------
    html_doc = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8"/>
  <title>Message Relationships Graph</title>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <script src="https://unpkg.com/cytoscape@3.28.1/dist/cytoscape.min.js"></script>
  <style>
    html, body {{ height: 100%; margin: 0; font-family: system-ui, -apple-system, Segoe UI, Roboto, sans-serif; }}
    #topbar {{
      padding: 10px; display: flex; gap: 8px; align-items: center;
      border-bottom: 1px solid #e5e7eb; flex-wrap: wrap;
    }}
    #cy {{ height: calc(100% - 56px); }}
    #stats-panel {{
      position: absolute; top: 70px; right: 10px; width: 300px; max-height: calc(100% - 80px);
      background: white; border: 1px solid #e5e7eb; border-radius: 8px; padding: 16px;
      box-shadow: 0 4px 6px -1px rgba(0,0,0,0.1); overflow-y: auto; z-index: 100;
      transition: transform 0.3s ease;
    }}
    #stats-panel.collapsed {{ transform: translateX(320px); }}
    #stats-toggle {{
      position: absolute; top: 70px; right: 10px; z-index: 101;
      background: white; border: 1px solid #e5e7eb; padding: 8px 12px;
      border-radius: 6px 0 0 6px; cursor: pointer; box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }}
    #stats-panel.collapsed + #stats-toggle {{ right: 10px; border-radius: 6px; }}
    .stats-section {{ margin-bottom: 16px; }}
    .stats-section h3 {{ margin: 0 0 8px 0; font-size: 14px; font-weight: 600; color: #374151; }}
    .stats-item {{ display: flex; justify-content: space-between; padding: 6px 0; border-bottom: 1px solid #f3f4f6; font-size: 13px; }}
    .stats-item:last-child {{ border-bottom: none; }}
    .stats-label {{ color: #6b7280; }}
    .stats-value {{ font-weight: 600; color: #111827; }}
    .top-relationship {{ padding: 6px; background: #f9fafb; margin-bottom: 4px; border-radius: 4px; font-size: 12px; }}
    .top-relationship .names {{ font-weight: 600; color: #374151; }}
    .top-relationship .count {{ color: #6b7280; font-size: 11px; }}
    #node-info {{
      position: absolute; top: 70px; left: 10px; width: 280px;
      background: white; border: 1px solid #e5e7eb; border-radius: 8px; padding: 16px;
      box-shadow: 0 4px 6px -1px rgba(0,0,0,0.1); z-index: 100;
      display: none;
    }}
    #node-info.visible {{ display: block; }}
    #node-info h3 {{ margin: 0 0 12px 0; font-size: 16px; font-weight: 600; color: #111827; }}
    #node-info .close-btn {{
      position: absolute; top: 12px; right: 12px; background: none; border: none;
      font-size: 20px; cursor: pointer; color: #9ca3af; padding: 0; line-height: 1;
    }}
    #node-info .close-btn:hover {{ color: #374151; }}
    #node-info .info-item {{ padding: 8px 0; border-bottom: 1px solid #f3f4f6; }}
    #node-info .info-item:last-child {{ border-bottom: none; }}
    #node-info .info-label {{ font-size: 12px; color: #6b7280; }}
    #node-info .info-value {{ font-size: 14px; font-weight: 600; color: #111827; margin-top: 2px; }}
    #node-info .connection-list {{ max-height: 150px; overflow-y: auto; }}
    #node-info .connection-item {{ padding: 4px 0; font-size: 12px; color: #374151; }}
    input, button, select {{ padding: 6px 8px; }}
    .pill {{ background:#eef2ff; color:#3730a3; padding:4px 8px; border-radius:999px; font-size:12px; }}
    .pill.warn {{ background:#fef3c7; color:#92400e; }}
    .tooltip {{ position: relative; display: inline-block; cursor: help; }}
    .tooltip .tooltiptext {{
      visibility: hidden; width: 280px; background-color: #1f2937; color: #fff;
      text-align: left; border-radius: 6px; padding: 8px; position: absolute;
      z-index: 1000; bottom: 125%; left: 50%; margin-left: -140px;
      opacity: 0; transition: opacity 0.3s; font-size: 12px; line-height: 1.4;
    }}
    .tooltip .tooltiptext::after {{
      content: ""; position: absolute; top: 100%; left: 50%; margin-left: -5px;
      border-width: 5px; border-style: solid; border-color: #1f2937 transparent transparent transparent;
    }}
    .tooltip:hover .tooltiptext {{ visibility: visible; opacity: 1; }}
    .info-icon {{
      display: inline-block; width: 14px; height: 14px; border-radius: 50%;
      background: #6b7280; color: white; text-align: center; font-size: 11px;
      line-height: 14px; font-weight: bold; margin-left: 3px;
    }}
  </style>
</head>
<body>
  <div id="topbar">
    <strong>Message Relationships</strong>
    <span class="pill">{total_nodes} contacts</span>
    <span class="pill">{total_edges} connections</span>
    <span class="pill">minimum messages: {min_edge}</span>
    <span class="pill">method 1: {inter_pairs} relationships</span>
    <span class="pill">method 2: {activity_pairs} relationships</span>
    {"<span class='pill warn'>calculation error</span>" if (activity_fn_name and act_err) else ""}
    <label>Graph Layout:
      <select id="layout">
        <option value="cose">Force-Directed</option>
        <option value="concentric">Concentric Circles</option>
        <option value="breadthfirst">Hierarchical</option>
        <option value="circle">Circle</option>
      </select>
    </label>
    <button id="fit">Reset View</button>
    <button id="toggle-edge-labels">Message Counts: On</button>
    <label>Counting Method:
      <select id="count-mode">
        <option value="merged">Total Messages (Maximum)</option>
        <option value="matrix">Total Messages (All)</option>
        <option value="activity">Conversational Turns</option>
      </select>
      <span class="tooltip">
        <span class="info-icon">?</span>
        <span class="tooltiptext">
          <strong>Total Messages (Maximum):</strong> Shows the highest count between the two methods - represents the actual total number of messages exchanged<br><br>
          <strong>Total Messages (All):</strong> Counts every message sent between each pair in their shared conversations - this is the forensically accurate total<br><br>
          <strong>Conversational Turns:</strong> Only counts messages where the speaker switches (e.g., A→B, B→A) - useful for analyzing conversation dynamics and responsiveness, but NOT the total message count. Multiple messages in a row from the same person count as one turn.
        </span>
      </span>
    </label>
    <label>Filter by Messages:
      <input id="min-edge-slider" type="range" min="1" max="{slider_max}" value="{min_edge}" style="width:100px;"/>
      <span id="min-edge-value">{min_edge}</span>
    </label>
    <input id="q" placeholder="Search by name…" />
  </div>
  <div id="cy"></div>

  <div id="stats-panel" class="collapsed">
    <div class="stats-section">
      <h3>Overall Statistics</h3>
      <div class="stats-item">
        <span class="stats-label">Total Relationships</span>
        <span class="stats-value">{stats_data['total_pairs']}</span>
      </div>
      <div class="stats-item">
        <span class="stats-label">Total Messages</span>
        <span class="stats-value">{stats_data['total_messages']}</span>
      </div>
      <div class="stats-item">
        <span class="stats-label">Most Active Contact</span>
        <span class="stats-value">{stats_data['most_active']}</span>
      </div>
    </div>

    <div class="stats-section">
      <h3>Top 10 Relationships</h3>
      {''.join([f'''
      <div class="top-relationship">
        <div class="names">{rel["person1"]} ↔ {rel["person2"]}</div>
        <div class="count">{rel["count"]} messages</div>
      </div>''' for rel in stats_data['top_relationships']])}
    </div>
  </div>

  <button id="stats-toggle">📊 Stats</button>

  <div id="node-info">
    <button class="close-btn" onclick="document.getElementById('node-info').classList.remove('visible')">×</button>
    <h3 id="node-name"></h3>
    <div class="info-item">
      <div class="info-label">Total Activity (All Messages)</div>
      <div class="info-value" id="node-total-activity-all"></div>
    </div>
    <div class="info-item">
      <div class="info-label">Visible in Current View</div>
      <div class="info-value" id="node-total-activity-visible"></div>
    </div>
    <div class="info-item">
      <div class="info-label">Number of Connections</div>
      <div class="info-value" id="node-connection-count"></div>
    </div>
    <div class="info-item">
      <div class="info-label">Top Connections</div>
      <div class="connection-list" id="node-connections"></div>
    </div>
  </div>

  <script>
    const elements = {json.dumps({"nodes": nodes, "edges": cy_edges})};

    function styleSpec() {{
      return [
        {{ selector: 'node',
          style: {{
            'background-color': 'data(color)',
            'label': 'data(label)',
            'width': 'data(size)',
            'height': 'data(size)',
            'text-valign': 'center',
            'color': '#111',
            'text-outline-color': '#fff',
            'text-outline-width': 2,
            'font-size': 12,
            'border-width': 2,
            'border-color': '#fff'
          }}
        }},
        {{ selector: 'edge',
          style: {{
            'width': 'data(width)',
            'line-color': '#cbd5e1',
            'curve-style': 'unbundled-bezier',
            'target-arrow-shape': 'triangle',
            'target-arrow-color': '#cbd5e1',
            'source-arrow-shape': 'triangle',
            'source-arrow-color': '#cbd5e1',
            'label': 'data(label)',
            'font-size': 9,
            'color': '#222',
            'text-rotation': 'autorotate',
            'text-background-color': '#ffffff',
            'text-background-opacity': 0.85,
            'text-background-padding': 2,
            'text-margin-y': -2,
            'text-wrap': 'wrap',
            'text-max-width': 80,
            'text-opacity': 1
          }}
        }},
        {{ selector: '.faded', style: {{ 'opacity': 0.15 }} }},
        {{ selector: '.highlight', style: {{ 'line-color': '#111', 'target-arrow-color': '#111', 'source-arrow-color': '#111', 'width': 6 }} }}
      ];
    }}

    function layoutSpec(name) {{
      if (name === 'breadthfirst') return {{ name, directed: false, padding: 20, spacingFactor: 1.15 }};
      if (name === 'concentric')  return {{ name, concentric: n => n.data('weight') || 1, levelWidth: () => 2, padding: 20 }};
      if (name === 'circle')      return {{ name, padding: 20 }};
      return {{ name: 'cose', animate: 'end', nodeRepulsion: 9000, idealEdgeLength: 120 }};
    }}

    const cy = cytoscape({{
      container: document.getElementById('cy'),
      elements,
      style: styleSpec(),
      layout: layoutSpec('cose'),
      wheelSensitivity: 0.2
    }});

    let focused = null;
    cy.on('tap', 'node', (evt) => {{
      const n = evt.target;

      // Show node info panel
      const nodeInfo = document.getElementById('node-info');
      document.getElementById('node-name').textContent = n.data('label');
      document.getElementById('node-total-activity-all').textContent = n.data('weight_all').toLocaleString() + ' messages';
      document.getElementById('node-total-activity-visible').textContent = n.data('weight').toLocaleString() + ' messages';

      // Get all connected edges and calculate connection details
      const connectedEdges = n.connectedEdges();
      document.getElementById('node-connection-count').textContent = connectedEdges.length;

      // Get current count mode
      const mode = document.getElementById('count-mode').value;

      // Build list of connections sorted by count
      const connections = [];
      connectedEdges.forEach(edge => {{
        const source = edge.data('source');
        const target = edge.data('target');
        const otherNode = source === n.id() ? target : source;
        let count;
        if (mode === 'matrix') {{
          count = edge.data('count_matrix');
        }} else if (mode === 'activity') {{
          count = edge.data('count_activity');
        }} else {{
          count = edge.data('count_merged');
        }}
        connections.push({{ name: otherNode, count: count }});
      }});

      connections.sort((a, b) => b.count - a.count);

      const connectionsList = document.getElementById('node-connections');
      connectionsList.innerHTML = connections.slice(0, 10).map(c =>
        `<div class="connection-item">${{c.name}}: ${{c.count}} messages</div>`
      ).join('');

      nodeInfo.classList.add('visible');

      // Highlight neighborhood
      if (focused && focused.id() === n.id()) {{
        cy.elements().removeClass('faded highlight');
        focused = null;
        nodeInfo.classList.remove('visible');
        return;
      }}
      focused = n;
      const k = n.closedNeighborhood();
      cy.elements().addClass('faded');
      k.removeClass('faded');
      k.edges().addClass('highlight');
    }});
    cy.on('tap', (evt) => {{
      if (evt.target === cy) {{
        cy.elements().removeClass('faded highlight');
        focused = null;
        document.getElementById('node-info').classList.remove('visible');
      }}
    }});

    document.getElementById('fit').addEventListener('click', () => cy.fit(cy.elements(), 40));
    document.getElementById('layout').addEventListener('change', (e) => {{
      cy.layout(layoutSpec(e.target.value)).run();
    }});

    const toggleBtn = document.getElementById('toggle-edge-labels');
    let labelsOn = true;
    toggleBtn.addEventListener('click', () => {{
      labelsOn = !labelsOn;
      cy.style().selector('edge').style('text-opacity', labelsOn ? 1 : 0).update();
      toggleBtn.textContent = `Message Counts: ${{labelsOn ? 'On' : 'Off'}}`;
    }});

    const countModeSelect = document.getElementById('count-mode');
    countModeSelect.addEventListener('change', (e) => {{
      const mode = e.target.value;
      cy.edges().forEach(edge => {{
        let newLabel;
        if (mode === 'matrix') {{
          newLabel = String(edge.data('count_matrix'));
        }} else if (mode === 'activity') {{
          newLabel = String(edge.data('count_activity'));
        }} else {{
          newLabel = String(edge.data('count_merged'));
        }}
        edge.data('label', newLabel);
      }});
    }});

    const minEdgeSlider = document.getElementById('min-edge-slider');
    const minEdgeValue = document.getElementById('min-edge-value');
    minEdgeSlider.addEventListener('input', (e) => {{
      const threshold = parseInt(e.target.value);
      minEdgeValue.textContent = threshold;

      // Get current count mode
      const mode = countModeSelect.value;

      // Hide/show edges based on threshold
      cy.edges().forEach(edge => {{
        let count;
        if (mode === 'matrix') {{
          count = edge.data('count_matrix');
        }} else if (mode === 'activity') {{
          count = edge.data('count_activity');
        }} else {{
          count = edge.data('count_merged');
        }}

        if (count >= threshold) {{
          edge.style('display', 'element');
        }} else {{
          edge.style('display', 'none');
        }}
      }});

      // Hide nodes that have no visible edges
      cy.nodes().forEach(node => {{
        const visibleEdges = node.connectedEdges().filter(e => e.style('display') === 'element');
        if (visibleEdges.length === 0) {{
          node.style('display', 'none');
        }} else {{
          node.style('display', 'element');
        }}
      }});
    }});

    const statsPanel = document.getElementById('stats-panel');
    const statsToggle = document.getElementById('stats-toggle');
    statsToggle.addEventListener('click', () => {{
      statsPanel.classList.toggle('collapsed');
      statsToggle.textContent = statsPanel.classList.contains('collapsed') ? '📊 Stats' : '✕ Close';
    }});

    const box = document.getElementById('q');
    box.addEventListener('input', () => {{
      const q = box.value.trim().toLowerCase();
      cy.elements().removeClass('faded highlight');
      if (!q) return;
      const hits = cy.nodes().filter(n => n.id().toLowerCase().includes(q) || String(n.data('label')||'').toLowerCase().includes(q));
      cy.elements().addClass('faded');
      hits.removeClass('faded');
      hits.neighborhood().removeClass('faded').edges().addClass('highlight');
      if (hits.length) cy.fit(hits, 80);
    }});
  </script>
</body>
</html>"""

    Path(out_html).write_text(html_doc, encoding="utf-8")
    return str(out_html)


def main():
    if len(sys.argv) == 1 or "--gui" in sys.argv[1:]:
        launch_gui()
        return
    parser = build_arg_parser()
    args = parser.parse_args()
    if args.gui:
        launch_gui()
        return
    run_cli(args)

if __name__ == "__main__":
    main()